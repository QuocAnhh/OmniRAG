#!/usr/bin/env python3
"""Build and benchmark the OpenDataLoader-backed OmniRAG pipeline.

The script uses generated fixtures and an isolated compose project by default:

  python scripts/benchmark_opendataloader_pipeline.py
  python scripts/benchmark_opendataloader_pipeline.py --skip-build --skip-up
"""

from __future__ import annotations

import argparse
import json
import mimetypes
import os
import subprocess
import sys
import time
import uuid
from pathlib import Path
from typing import Any

try:
    import requests
except ImportError as exc:  # pragma: no cover - operator-facing utility
    raise SystemExit("The benchmark needs requests: python -m pip install requests") from exc


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PROJECT = "omnirag-odl-bench"
DEFAULT_BASE_URL = "http://localhost:18080"
DEFAULT_FIXTURES_DIR = Path("/tmp/omnirag-odl-fixtures")
DEFAULT_OUT_JSON = Path("/tmp/omnirag-odl-benchmark.json")
DEFAULT_OUT_MD = Path("/tmp/omnirag-odl-benchmark.md")
COMPOSE_SERVICES = [
    "db",
    "mongodb",
    "redis",
    "minio",
    "qdrant",
    "opendataloader-hybrid",
    "backend",
    "celery_worker",
    "gateway",
    "frontend",
]
DEFAULT_PORT_ENV = {
    "GATEWAY_HOST_PORT": "18080",
    "BACKEND_HOST_PORT": "18001",
    "QDRANT_HOST_PORT": "16333",
    "PDF_HYBRID_HOST_PORT": "15002",
    "MINIO_API_HOST_PORT": "19000",
    "MINIO_CONSOLE_HOST_PORT": "19001",
    "POSTGRES_HOST_PORT": "15433",
    "REDIS_HOST_PORT": "16380",
    "MONGODB_HOST_PORT": "27018",
    "FRONTEND_HOST_PORT": "15173",
}


def run_cmd(cmd: list[str], env: dict[str, str] | None = None, timeout: int | None = None) -> float:
    start = time.perf_counter()
    subprocess.run(cmd, cwd=ROOT, env=env, timeout=timeout, check=True)
    return time.perf_counter() - start


def load_env_file(path: Path) -> dict[str, str]:
    values: dict[str, str] = {}
    if not path.exists():
        return values
    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        values[key.strip()] = value.strip().strip('"').strip("'")
    return values


def has_openrouter_key() -> bool:
    env_values = load_env_file(ROOT / "backend" / ".env")
    return bool(os.getenv("OPENROUTER_API_KEY") or env_values.get("OPENROUTER_API_KEY"))


def qdrant_collection_name() -> str:
    env_values = load_env_file(ROOT / "backend" / ".env")
    return (
        os.getenv("RAG_COLLECTION_NAME")
        or env_values.get("RAG_COLLECTION_NAME")
        or "omnirag_openrouter_collection_v3"
    )


def write_simple_pdf(path: Path) -> None:
    text = (
        "OmniRAG OpenDataLoader Benchmark|"
        "Heading: Revenue Report 2026|"
        "Table: Quarter Revenue Cost|Q1 100 42|Q2 135 51|"
        "Figure 1: Bar chart showing revenue growth"
    )
    stream = (
        "BT\n/F1 18 Tf\n72 720 Td\n"
        + "\\n".join(f"({line}) Tj 0 -28 Td" for line in text.split("|"))
        + "\nET"
    ).encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    content = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for idx, obj in enumerate(objects, start=1):
        offsets.append(len(content))
        content.extend(f"{idx} 0 obj\n".encode("ascii"))
        content.extend(obj)
        content.extend(b"\nendobj\n")
    xref = len(content)
    content.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    content.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    content.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii")
    )
    path.write_bytes(bytes(content))


def create_docx(path: Path) -> bool:
    try:
        from docx import Document
    except ImportError:
        return False
    doc = Document()
    doc.add_heading("OmniRAG DOCX Fixture", level=1)
    doc.add_paragraph("The integration should extract this paragraph.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Latency"
    table.cell(1, 1).text = "Tracked"
    doc.save(path)
    return True


def create_pptx(path: Path) -> bool:
    try:
        from pptx import Presentation
    except ImportError:
        return False
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "OmniRAG PPTX Fixture"
    slide.placeholders[1].text = "Slides should be parsed into searchable text."
    prs.save(path)
    return True


def create_xlsx(path: Path) -> bool:
    try:
        from openpyxl import Workbook
    except ImportError:
        return False
    wb = Workbook()
    ws = wb.active
    ws.title = "Metrics"
    ws.append(["Metric", "Value"])
    ws.append(["Ingest latency", "measured"])
    ws.append(["Structured metadata", "enabled"])
    wb.save(path)
    return True


def create_fixtures(fixtures_dir: Path) -> dict[str, Any]:
    fixtures_dir.mkdir(parents=True, exist_ok=True)
    created: list[str] = []
    skipped: list[str] = []

    always = {
        "benchmark.md": "# OmniRAG Markdown Fixture\n\nStructured markdown should become searchable.",
        "benchmark.txt": "OmniRAG text fixture for retrieval latency measurement.",
        "benchmark.csv": "metric,value\nretrieval,hybrid\ncitation,page_bbox\n",
    }
    for filename, content in always.items():
        path = fixtures_dir / filename
        path.write_text(content, encoding="utf-8")
        created.append(str(path))

    pdf_path = fixtures_dir / "benchmark.pdf"
    write_simple_pdf(pdf_path)
    created.append(str(pdf_path))

    optional_creators = {
        "benchmark.docx": create_docx,
        "benchmark.pptx": create_pptx,
        "benchmark.xlsx": create_xlsx,
    }
    for filename, creator in optional_creators.items():
        path = fixtures_dir / filename
        if creator(path):
            created.append(str(path))
        else:
            skipped.append(filename)

    legacy_path = fixtures_dir / "legacy.doc"
    legacy_path.write_text("legacy binary placeholder", encoding="utf-8")

    return {
        "created": created,
        "skipped_optional": skipped,
        "legacy": str(legacy_path),
    }


def wait_http(url: str, timeout_s: int, session: requests.Session) -> float:
    start = time.perf_counter()
    last_error = ""
    while time.perf_counter() - start < timeout_s:
        try:
            resp = session.get(url, timeout=5)
            if 200 <= resp.status_code < 300:
                return time.perf_counter() - start
            last_error = f"{resp.status_code}: {resp.text[:200]}"
        except Exception as exc:
            last_error = str(exc)
        time.sleep(3)
    raise RuntimeError(f"Timed out waiting for {url}: {last_error}")


def post_json(session: requests.Session, url: str, payload: dict[str, Any], timeout: int = 60) -> requests.Response:
    resp = session.post(url, json=payload, timeout=timeout)
    if resp.status_code >= 400:
        raise RuntimeError(f"POST {url} failed {resp.status_code}: {resp.text[:500]}")
    return resp


def register_login_create_bot(base_url: str, session: requests.Session) -> dict[str, Any]:
    suffix = uuid.uuid4().hex[:10]
    email = f"odl-bench-{suffix}@example.com"
    password = "OmniRAG-bench-12345"

    register_payload = {
        "email": email,
        "password": password,
        "full_name": "OpenDataLoader Benchmark",
        "tenant_name": f"ODL Bench {suffix}",
    }
    reg_start = time.perf_counter()
    register_resp = session.post(f"{base_url}/api/v1/auth/register", json=register_payload, timeout=60)
    if register_resp.status_code >= 400:
        raise RuntimeError(f"Register failed {register_resp.status_code}: {register_resp.text[:500]}")
    register_latency = time.perf_counter() - reg_start

    login_start = time.perf_counter()
    login_resp = session.post(
        f"{base_url}/api/v1/auth/login",
        data={"username": email, "password": password},
        timeout=60,
    )
    if login_resp.status_code >= 400:
        raise RuntimeError(f"Login failed {login_resp.status_code}: {login_resp.text[:500]}")
    login_latency = time.perf_counter() - login_start
    token = login_resp.json()["access_token"]
    session.headers.update({"Authorization": f"Bearer {token}"})

    bot_payload = {
        "name": f"ODL Bench {suffix}",
        "description": "Generated OpenDataLoader benchmark bot",
        "config": {
            "model": "openai/gpt-4o-mini",
            "domain": "general",
            "enable_memory": False,
            "top_k": 5,
            "similarity_threshold": 0,
            "chunking_strategy": "recursive",
            "chunk_size": 800,
            "chunk_overlap": 120,
            "pdf_parser_mode": "hybrid_auto",
            "pdf_structured_chunking": True,
            "pdf_enrich_formula": False,
            "pdf_sanitize": False,
            "pdf_use_struct_tree": False,
            "pdf_include_header_footer": False,
            "pdf_detect_strikethrough": False,
            "pdf_threads": 1,
            "enrich_picture_description": False,
        },
    }
    bot_start = time.perf_counter()
    bot_resp = post_json(session, f"{base_url}/api/v1/bots/", bot_payload)
    bot_latency = time.perf_counter() - bot_start

    return {
        "email": email,
        "bot": bot_resp.json(),
        "latencies_s": {
            "register": round(register_latency, 3),
            "login": round(login_latency, 3),
            "create_bot": round(bot_latency, 3),
        },
    }


def upload_file(session: requests.Session, base_url: str, bot_id: str, path: Path) -> tuple[dict[str, Any], float]:
    content_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    with path.open("rb") as fh:
        start = time.perf_counter()
        resp = session.post(
            f"{base_url}/api/v1/bots/{bot_id}/documents",
            files={"file": (path.name, fh, content_type)},
            data={"chunking_strategy": "recursive", "enable_knowledge_graph": "false"},
            timeout=120,
        )
        latency = time.perf_counter() - start
    if resp.status_code >= 400:
        raise RuntimeError(f"Upload {path.name} failed {resp.status_code}: {resp.text[:500]}")
    return resp.json(), latency


def poll_documents(
    session: requests.Session,
    base_url: str,
    bot_id: str,
    uploaded_ids: set[str],
    timeout_s: int,
) -> tuple[list[dict[str, Any]], float]:
    start = time.perf_counter()
    while time.perf_counter() - start < timeout_s:
        resp = session.get(f"{base_url}/api/v1/bots/{bot_id}/documents", timeout=30)
        if resp.status_code >= 400:
            raise RuntimeError(f"List documents failed {resp.status_code}: {resp.text[:500]}")
        docs = [doc for doc in resp.json() if doc["id"] in uploaded_ids]
        if docs and all(doc["status"] in {"completed", "failed"} for doc in docs):
            return docs, time.perf_counter() - start
        time.sleep(5)
    raise RuntimeError(f"Timed out waiting for document ingestion after {timeout_s}s")


def qdrant_count(bot_id: str, collection: str, port: str) -> int:
    payload = {
        "filter": {"must": [{"key": "bot_id", "match": {"value": bot_id}}]},
        "exact": True,
    }
    resp = requests.post(
        f"http://localhost:{port}/collections/{collection}/points/count",
        json=payload,
        timeout=30,
    )
    if resp.status_code >= 400:
        raise RuntimeError(f"Qdrant count failed {resp.status_code}: {resp.text[:500]}")
    return int(resp.json().get("result", {}).get("count", 0))


def run_retrieval_chat(session: requests.Session, base_url: str, bot_id: str) -> dict[str, Any]:
    query = "What does the benchmark document say about revenue and structured metadata?"
    retrieval_start = time.perf_counter()
    retrieval_resp = post_json(
        session,
        f"{base_url}/api/v1/bots/{bot_id}/retrieve",
        {"query": query, "top_k": 5},
        timeout=120,
    )
    retrieval_latency = time.perf_counter() - retrieval_start
    retrieval_json = retrieval_resp.json()

    chat_start = time.perf_counter()
    chat_resp = post_json(
        session,
        f"{base_url}/api/v1/bots/{bot_id}/chat",
        {"message": query, "history": []},
        timeout=180,
    )
    chat_latency = time.perf_counter() - chat_start
    chat_json = chat_resp.json()

    ttft: float | None = None
    stream_chunks = 0
    stream_start = time.perf_counter()
    with session.post(
        f"{base_url}/api/v1/bots/{bot_id}/chat-stream",
        json={"message": query, "history": []},
        stream=True,
        timeout=180,
    ) as resp:
        if resp.status_code >= 400:
            raise RuntimeError(f"Stream chat failed {resp.status_code}: {resp.text[:500]}")
        for line in resp.iter_lines(decode_unicode=True):
            if not line:
                continue
            if line.startswith("data:"):
                stream_chunks += 1
                if ttft is None:
                    ttft = time.perf_counter() - stream_start
                if '"type": "done"' in line or '"type":"done"' in line:
                    break
    stream_total = time.perf_counter() - stream_start

    return {
        "query": query,
        "retrieval_latency_s": round(retrieval_latency, 3),
        "retrieval_results": len(retrieval_json.get("results", [])),
        "chat_total_latency_s": round(chat_latency, 3),
        "chat_retrieved_chunks": len(chat_json.get("retrieved_chunks", [])),
        "chat_response_chars": len(chat_json.get("response", "")),
        "stream_ttft_s": round(ttft, 3) if ttft is not None else None,
        "stream_total_s": round(stream_total, 3),
        "stream_events": stream_chunks,
    }


def test_legacy_block(session: requests.Session, base_url: str, bot_id: str, legacy_path: Path) -> dict[str, Any]:
    with legacy_path.open("rb") as fh:
        resp = session.post(
            f"{base_url}/api/v1/bots/{bot_id}/documents",
            files={"file": (legacy_path.name, fh, "application/msword")},
            data={"chunking_strategy": "recursive", "enable_knowledge_graph": "false"},
            timeout=60,
        )
    return {
        "status_code": resp.status_code,
        "blocked": resp.status_code == 415,
        "detail": resp.text[:300],
    }


def write_reports(report: dict[str, Any], out_json: Path, out_md: Path) -> None:
    out_json.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")

    lines = [
        "# OmniRAG OpenDataLoader Benchmark",
        "",
        f"- status: `{report['status']}`",
        f"- compose_project: `{report['compose_project']}`",
        f"- base_url: `{report['base_url']}`",
        f"- collection: `{report['collection']}`",
        f"- generated_fixtures: `{len(report['fixtures']['created'])}`",
    ]
    if report.get("docker"):
        for key, value in report["docker"].items():
            lines.append(f"- docker_{key}: `{value}`")
    if report.get("pipeline"):
        pipeline = report["pipeline"]
        lines.extend([
            "",
            "## Pipeline",
            f"- uploads: `{len(pipeline.get('uploads', []))}`",
            f"- ingest_poll_s: `{pipeline.get('ingest_poll_s')}`",
            f"- qdrant_points: `{pipeline.get('qdrant_points')}`",
            f"- total_chunks: `{pipeline.get('total_chunks')}`",
        ])
    if report.get("chat"):
        chat = report["chat"]
        lines.extend([
            "",
            "## Chat",
            f"- retrieval_latency_s: `{chat.get('retrieval_latency_s')}`",
            f"- retrieval_results: `{chat.get('retrieval_results')}`",
            f"- chat_total_latency_s: `{chat.get('chat_total_latency_s')}`",
            f"- chat_retrieved_chunks: `{chat.get('chat_retrieved_chunks')}`",
            f"- stream_ttft_s: `{chat.get('stream_ttft_s')}`",
        ])
    if report.get("notes"):
        lines.extend(["", "## Notes", *[f"- {note}" for note in report["notes"]]])
    if report.get("errors"):
        lines.extend(["", "## Errors", *[f"- {err}" for err in report["errors"]]])
    out_md.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> int:
    parser = argparse.ArgumentParser(description="Build and benchmark OmniRAG OpenDataLoader pipeline")
    parser.add_argument("--project-name", default=DEFAULT_PROJECT)
    parser.add_argument("--base-url", default=DEFAULT_BASE_URL)
    parser.add_argument("--fixtures-dir", type=Path, default=DEFAULT_FIXTURES_DIR)
    parser.add_argument("--out-json", type=Path, default=DEFAULT_OUT_JSON)
    parser.add_argument("--out-md", type=Path, default=DEFAULT_OUT_MD)
    parser.add_argument("--skip-build", action="store_true")
    parser.add_argument("--skip-up", action="store_true")
    parser.add_argument("--timeout", type=int, default=900, help="Ingest polling timeout in seconds")
    args = parser.parse_args()

    report: dict[str, Any] = {
        "status": "running",
        "compose_project": args.project_name,
        "base_url": args.base_url.rstrip("/"),
        "collection": qdrant_collection_name(),
        "fixtures": {},
        "docker": {},
        "pipeline": {},
        "chat": {},
        "notes": [],
        "errors": [],
        "started_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
    }

    session = requests.Session()
    compose_env = os.environ.copy()
    compose_env["COMPOSE_PROJECT_NAME"] = args.project_name
    for key, value in load_env_file(ROOT / "backend" / ".env").items():
        compose_env.setdefault(key, value)
    for key, value in DEFAULT_PORT_ENV.items():
        compose_env.setdefault(key, value)
    compose_env.setdefault("PDF_HYBRID_DEVICE", "cpu")

    try:
        report["fixtures"] = create_fixtures(args.fixtures_dir)

        if not args.skip_build:
            elapsed = run_cmd(
                ["docker", "compose", "-p", args.project_name, "build", "backend", "opendataloader-hybrid", "gateway", "frontend"],
                env=compose_env,
                timeout=3600,
            )
            report["docker"]["build_s"] = round(elapsed, 3)

        if not args.skip_up:
            elapsed = run_cmd(
                ["docker", "compose", "-p", args.project_name, "up", "-d", *COMPOSE_SERVICES],
                env=compose_env,
                timeout=1200,
            )
            report["docker"]["up_s"] = round(elapsed, 3)

        base_url = report["base_url"]
        report["docker"]["host_ports"] = {key: compose_env[key] for key in DEFAULT_PORT_ENV}
        report["docker"]["gateway_health_s"] = round(wait_http(f"{base_url}/health", 300, session), 3)
        report["docker"]["backend_health_s"] = round(
            wait_http(f"http://localhost:{compose_env['BACKEND_HOST_PORT']}/health", 300, session), 3
        )
        report["docker"]["qdrant_health_s"] = round(
            wait_http(f"http://localhost:{compose_env['QDRANT_HOST_PORT']}/healthz", 120, session), 3
        )
        report["docker"]["minio_health_s"] = round(
            wait_http(f"http://localhost:{compose_env['MINIO_API_HOST_PORT']}/minio/health/live", 120, session), 3
        )
        report["docker"]["hybrid_health_s"] = round(
            wait_http(f"http://localhost:{compose_env['PDF_HYBRID_HOST_PORT']}/health", 300, session), 3
        )

        account = register_login_create_bot(base_url, session)
        bot_id = account["bot"]["id"]
        report["pipeline"]["auth_bot"] = account["latencies_s"]
        report["pipeline"]["bot_id"] = bot_id

        legacy = test_legacy_block(session, base_url, bot_id, Path(report["fixtures"]["legacy"]))
        report["pipeline"]["legacy_block"] = legacy
        if not legacy["blocked"]:
            raise RuntimeError(f"Legacy .doc upload was not blocked: {legacy}")

        if not has_openrouter_key():
            report["status"] = "skipped_no_openrouter_key"
            report["notes"].append("OPENROUTER_API_KEY is missing; skipped ingest, Qdrant retrieval, and chat.")
            return 0

        uploads = []
        uploaded_ids: set[str] = set()
        for fixture in report["fixtures"]["created"]:
            path = Path(fixture)
            doc, latency = upload_file(session, base_url, bot_id, path)
            uploaded_ids.add(doc["id"])
            uploads.append({"filename": path.name, "document_id": doc["id"], "upload_latency_s": round(latency, 3)})
        report["pipeline"]["uploads"] = uploads

        docs, ingest_latency = poll_documents(session, base_url, bot_id, uploaded_ids, args.timeout)
        report["pipeline"]["ingest_poll_s"] = round(ingest_latency, 3)
        report["pipeline"]["documents"] = [
            {
                "filename": doc["filename"],
                "status": doc["status"],
                "num_chunks": (doc.get("doc_metadata") or {}).get("num_chunks", 0),
                "has_structured_json": (doc.get("doc_metadata") or {}).get("has_structured_json"),
                "page_numbers": (doc.get("doc_metadata") or {}).get("page_numbers", []),
                "error_message": doc.get("error_message"),
            }
            for doc in docs
        ]

        failed_docs = [doc for doc in docs if doc["status"] == "failed"]
        if failed_docs:
            raise RuntimeError(f"Document ingest failed: {failed_docs}")
        total_chunks = sum(int((doc.get("doc_metadata") or {}).get("num_chunks") or 0) for doc in docs)
        report["pipeline"]["total_chunks"] = total_chunks
        if total_chunks <= 0:
            raise RuntimeError("Ingest completed with zero chunks")

        count = qdrant_count(bot_id, report["collection"], compose_env["QDRANT_HOST_PORT"])
        report["pipeline"]["qdrant_points"] = count
        if count <= 0:
            raise RuntimeError("Qdrant count is zero after ingest")

        report["chat"] = run_retrieval_chat(session, base_url, bot_id)
        if report["chat"]["retrieval_results"] <= 0:
            raise RuntimeError("Retrieval returned zero results")
        if report["chat"]["chat_retrieved_chunks"] <= 0:
            raise RuntimeError("Chat returned zero retrieved chunks")

        report["status"] = "passed"
        return 0
    except Exception as exc:
        report["status"] = "failed"
        report["errors"].append(str(exc))
        return 1
    finally:
        report["finished_at"] = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
        write_reports(report, args.out_json, args.out_md)
        print(json.dumps({"status": report["status"], "json": str(args.out_json), "markdown": str(args.out_md)}, indent=2))


if __name__ == "__main__":
    raise SystemExit(main())
