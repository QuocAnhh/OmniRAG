"""
OpenRouter-Enhanced RAG Service
Integrates OpenRouter for both LLM and embeddings in RAG pipeline.
Provides seamless migration path from existing services.
"""

import os
import logging

logger = logging.getLogger(__name__)

# --- ENABLE WRITABLE CACHE FOR DOCKER ---
model_cache_dir = os.getenv('HF_HOME', '/tmp/huggingface_cache')
try:
    # Try to create and check if writable
    os.makedirs(model_cache_dir, exist_ok=True)
    # Test writability by creating a temp file
    test_file = os.path.join(model_cache_dir, '.write_test')
    with open(test_file, 'w') as f:
        f.write('test')
    os.remove(test_file)
except (Exception, OSError) as e:
    # Fallback to /tmp if primary fails (usually due to Docker volume permissions)
    logger.warning("primary_cache_not_writable", extra={"cache_dir": model_cache_dir, "error": str(e)})
    model_cache_dir = '/tmp/huggingface_cache'
    os.makedirs(model_cache_dir, exist_ok=True)

# Set all possible cache variables
os.environ['HF_HOME'] = model_cache_dir
os.environ['SENTENCE_TRANSFORMERS_HOME'] = model_cache_dir
os.environ['TRANSFORMERS_CACHE'] = model_cache_dir
os.environ['XDG_CACHE_HOME'] = model_cache_dir

try:
    # Clean up stale lock files if they exist
    if os.path.exists(model_cache_dir):
        for root, dirs, files in os.walk(model_cache_dir):
            for file in files:
                if file.endswith(".lock"):
                    try:
                        os.remove(os.path.join(root, file))
                    except OSError:
                        pass
except Exception as e:
    logger.warning("lock_cleanup_failed", extra={"error": str(e)})

import time
import uuid
import asyncio
import csv
from typing import List, Dict, Any, Optional, AsyncGenerator, Union
from datetime import datetime
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
import httpx

from fastapi import UploadFile
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LangChainDocument
from langchain_community.document_loaders import TextLoader
from qdrant_client import QdrantClient, AsyncQdrantClient
from qdrant_client.http import models as rest
from qdrant_client.models import PointStruct, Filter, FieldCondition, MatchValue

from app.core.config import settings
from app.db.mongodb import get_mongodb
from app.services.openrouter_service import get_openrouter_service
from app.services.memory_service import memory_service
from pathlib import Path
import tempfile
import shutil
import hashlib
import json
import re

# Fast model for internal pipeline calls (query rewriting, HyDE, CRAG, etc.)
# User-facing answer generation uses the model configured per-bot in the UI.
INTERNAL_LLM_MODEL = "google/gemini-3.1-flash-lite"


def _sigmoid(x):
    """Sigmoid normalization: maps logit scores to 0..1 probability range."""
    import numpy as np
    return 1 / (1 + np.exp(-x))


def _build_group_context_block(bot_config: dict) -> str:
    """Build the [NHÓM CHAT] system-prompt block from Facebook group context keys."""
    lines: list[str] = []

    group_name = (bot_config.get("group_name") or "").strip()
    description = (bot_config.get("group_description") or "").strip()
    participants: list[dict] = bot_config.get("group_participants") or []
    recent_msgs: list[dict] = bot_config.get("group_recent_messages") or []

    header = f'Nhóm: "{group_name}"' if group_name else "Nhóm chat"
    if description:
        header += f" — {description}"
    lines.append(f"[{header}")

    if participants:
        lines.append(f"Thành viên ({len(participants)} người):")
        for p in participants:
            name = p.get("name") or ""
            nick = p.get("nickname") or ""
            admin = p.get("is_admin", False)
            tags: list[str] = []
            if admin:
                tags.append("admin")
            if nick and nick != name:
                tags.append(f'biệt danh: "{nick}"')
            suffix = f" ({', '.join(tags)})" if tags else ""
            lines.append(f"  • @{name}{suffix}")

    if recent_msgs:
        lines.append(f"\nLịch sử trò chuyện gần đây ({len(recent_msgs)} tin nhắn):")
        for m in recent_msgs:
            sender = m.get("sender_name") or m.get("sender_id") or "?"
            text = (m.get("text") or "").strip()
            if text:
                lines.append(f"  [{sender}]: {text}")

    lines.append(
        "\nQUY TẮC TAG:\n"
        "• Khi được yêu cầu làm gì đó với ai (tán, hỏi, nhắn tin, tag,...) → tag THẲNG người đó bằng @Tên, làm luôn, KHÔNG hỏi lại.\n"
        "• KHÔNG tự động tag lại người đang nói chuyện với mày trừ khi câu trả lời đang hướng thẳng vào họ.\n"
        "• Nếu không tìm thấy tên trong danh sách → báo luôn 'tao không thấy ai tên X trong nhóm'.\n"
        "• Hệ thống tự chuyển @Tên thành mention thực trên Facebook.]"
    )
    return "\n\n" + "\n".join(lines)


def _extract_lightrag_entity_names(lightrag_raw: str) -> list:
    """
    Parse entity names from LightRAG's only_need_context=True output.

    Actual format produced by LightRAG:

        Knowledge Graph Data (Entity):

        ```json
        {"entity": "Prompt Engineering", "type": "concept", "description": "..."}
        {"entity": "Generative AI", "type": "concept", "description": "..."}
        ```

        Knowledge Graph Data (Relationship):
        ...

    Each entity line is a standalone JSON object (not an array).
    """
    entities = []
    try:
        in_entity_block = False
        in_code_block = False
        for line in lightrag_raw.splitlines():
            stripped = line.strip()

            # Detect section headers
            if "Knowledge Graph Data (Entity)" in stripped:
                in_entity_block = True
                in_code_block = False
                continue
            if in_entity_block and "Knowledge Graph Data (Relationship)" in stripped:
                break  # done with entities section

            if not in_entity_block:
                continue

            # Track ```json ... ``` fences
            if stripped.startswith("```"):
                in_code_block = not in_code_block
                continue

            if not in_code_block or not stripped.startswith("{"):
                continue

            # Parse JSON line
            try:
                obj = json.loads(stripped)
                name = obj.get("entity", "").strip()
                # Skip noise: empty, pure numbers, single chars, type tags, "UNKNOWN"
                if name and len(name) > 1 and name.upper() != "UNKNOWN" and len(name.split()) <= 5:
                    entities.append(name)
            except (json.JSONDecodeError, ValueError):
                # Try extracting "entity" value with a targeted regex as fallback
                m = re.search(r'"entity"\s*:\s*"([^"]+)"', stripped)
                if m:
                    name = m.group(1).strip()
                    if name and len(name) > 1 and name.upper() != "UNKNOWN" and len(name.split()) <= 5:
                        entities.append(name)

    except Exception:
        logger.warning("lightrag_entity_parse_failed")

    # Deduplicate, preserve order
    seen: set = set()
    result = []
    for e in entities:
        key = e.lower()
        if key not in seen:
            seen.add(key)
            result.append(e)

    return result[:20]  # cap at 20 — enough for meaningful highlighting


class OpenRouterAPIError(Exception):
    """Custom exception for OpenRouter API errors"""
    pass



# from sentence_transformers import CrossEncoder # Moved to lazy loading in _get_reranker

class OpenRouterRAGService:
    """
    RAG Service powered by OpenRouter for both LLM and embeddings.
    Drop-in replacement for existing RAG services with enhanced capabilities.
    """
    
    def __init__(self):
        """Initialize OpenRouter RAG service with all dependencies."""
        self.openrouter = get_openrouter_service()
        self.qdrant_client = QdrantClient(host=settings.QDRANT_HOST, port=settings.QDRANT_PORT)
        self.async_qdrant_client = AsyncQdrantClient(host=settings.QDRANT_HOST, port=settings.QDRANT_PORT)
        self.collection_name = settings.RAG_COLLECTION_NAME
        self.legacy_collection_name = settings.RAG_LEGACY_COLLECTION_NAME
        self.dense_vector_name = "dense"
        self.sparse_vector_name = "bm25"
        
        # Redis caching via async client (shared from app.db.redis)
        from app.services.cache_service import cache_service
        self.cache = cache_service
        
        # Get embedding dimensions (typically 1536 for text-embedding-3-small)
        self.embedding_dim = 1536  # Will be auto-detected on first embedding
        self._sparse_embedder = None
        
        self._ensure_collection()
        
        # Deferred reranker initialization (lazy load)
        self.reranker = None
        self._reranker_attempted = False

        logger.info("OpenRouterRAGService initialized successfully (reranker loads on demand)")
    
    def _ensure_collection(self):
        """Ensure Qdrant collection exists with dense+sparse named vectors."""
        try:
            self.qdrant_client.get_collection(self.collection_name)
            logger.info(f"Collection '{self.collection_name}' already exists")
        except Exception:
            logger.info(f"Creating collection '{self.collection_name}' with dense+sparse vectors")

            self.qdrant_client.create_collection(
                collection_name=self.collection_name,
                vectors_config={
                    self.dense_vector_name: rest.VectorParams(
                        size=self.embedding_dim,
                        distance=rest.Distance.COSINE,
                    )
                },
                sparse_vectors_config={
                    self.sparse_vector_name: rest.SparseVectorParams(
                        modifier=rest.Modifier.IDF
                    )
                },
                quantization_config=rest.ScalarQuantization(
                    scalar=rest.ScalarQuantizationConfig(
                        type=rest.ScalarType.INT8,
                        quantile=0.99,
                        always_ram=True
                    )
                )
            )
            logger.info(f"Collection '{self.collection_name}' created successfully")

        self._ensure_payload_indexes()

    def _ensure_payload_indexes(self):
        """Create payload indexes if missing; tolerate already-existing indexes."""
        for field_name in ("bot_id", "document_id", "source"):
            try:
                self.qdrant_client.create_payload_index(
                    collection_name=self.collection_name,
                    field_name=field_name,
                    field_schema=rest.PayloadSchemaType.KEYWORD,
                )
            except Exception:
                logger.debug("qdrant_payload_index_exists_or_failed", extra={"field": field_name})
        try:
            self.qdrant_client.create_payload_index(
                collection_name=self.collection_name,
                field_name="text",
                field_schema=rest.TextIndexParams(
                    type="text",
                    tokenizer=rest.TokenizerType.MULTILINGUAL,
                    lowercase=True,
                    min_token_len=2,
                    max_token_len=20,
                ),
            )
        except Exception:
            logger.debug("qdrant_text_index_exists_or_failed")

    def _get_sparse_embedder(self):
        """Lazy BM25 sparse embedder used for Qdrant hybrid search."""
        if self._sparse_embedder is not None:
            return self._sparse_embedder
        try:
            from fastembed import SparseTextEmbedding

            kwargs = {
                "model_name": settings.QDRANT_SPARSE_MODEL,
                "language": settings.QDRANT_SPARSE_LANGUAGE,
                "disable_stemmer": settings.QDRANT_SPARSE_DISABLE_STEMMER,
            }
            try:
                self._sparse_embedder = SparseTextEmbedding(**kwargs)
            except TypeError:
                self._sparse_embedder = SparseTextEmbedding(model_name=settings.QDRANT_SPARSE_MODEL)
            logger.info("sparse_embedder_loaded", extra={"model": settings.QDRANT_SPARSE_MODEL})
            return self._sparse_embedder
        except Exception as e:
            logger.error("sparse_embedder_load_failed", extra={"error": str(e)})
            raise

    @staticmethod
    def _to_qdrant_sparse_vector(sparse_embedding):
        """Convert FastEmbed SparseEmbedding to Qdrant SparseVector."""
        indices = getattr(sparse_embedding, "indices", [])
        values = getattr(sparse_embedding, "values", [])
        if hasattr(indices, "tolist"):
            indices = indices.tolist()
        if hasattr(values, "tolist"):
            values = values.tolist()
        return rest.SparseVector(
            indices=[int(i) for i in indices],
            values=[float(v) for v in values],
        )

    def _embed_sparse_texts(self, texts: List[str]) -> List[Any]:
        if not texts:
            return []
        embedder = self._get_sparse_embedder()
        return [
            self._to_qdrant_sparse_vector(embedding)
            for embedding in embedder.embed(texts)
        ]

    def _embed_sparse_query(self, query: str) -> Any:
        vectors = self._embed_sparse_texts([query])
        return vectors[0] if vectors else rest.SparseVector(indices=[], values=[])

    @staticmethod
    def _build_named_vector(dense: List[float], sparse: Any) -> Dict[str, Any]:
        return {
            "dense": dense,
            "bm25": sparse,
        }

    @staticmethod
    def _build_rrf_query():
        """Use modern RRF query shape, with FusionQuery fallback for older clients."""
        if hasattr(rest, "FusionQuery") and hasattr(rest, "Fusion"):
            return rest.FusionQuery(fusion=rest.Fusion.RRF)
        return rest.RrfQuery(rrf=rest.Rrf())
    
    @staticmethod
    def _config_bool(config: Optional[dict], key: str, default: bool = False) -> bool:
        value = (config or {}).get(key, default)
        if isinstance(value, str):
            return value.strip().lower() in {"1", "true", "yes", "on"}
        return bool(value)

    @staticmethod
    def _config_int(config: Optional[dict], key: str, default: int) -> int:
        try:
            return int((config or {}).get(key, default))
        except (TypeError, ValueError):
            return default

    def _resolve_pdf_runtime_config(
        self,
        bot_config: Optional[dict] = None,
        enrich_picture_description: bool = False,
    ) -> Dict[str, Any]:
        """Resolve per-bot OpenDataLoader settings with conservative fallbacks."""
        bot_config = bot_config or {}
        parser_mode = bot_config.get("pdf_parser_mode") or settings.PDF_DEFAULT_PARSER_MODE
        if parser_mode not in {"local_fast", "hybrid_auto", "hybrid_full"}:
            parser_mode = settings.PDF_DEFAULT_PARSER_MODE

        picture = enrich_picture_description or self._config_bool(
            bot_config, "enrich_picture_description", False
        )
        formula = self._config_bool(bot_config, "pdf_enrich_formula", settings.PDF_ENRICH_FORMULA)
        if picture or formula:
            parser_mode = "hybrid_full"

        return {
            "parser_mode": parser_mode,
            "structured_chunking": self._config_bool(
                bot_config, "pdf_structured_chunking", settings.PDF_STRUCTURED_CHUNKING
            ),
            "enrich_picture_description": picture,
            "enrich_formula": formula,
            "sanitize": self._config_bool(bot_config, "pdf_sanitize", settings.PDF_SANITIZE),
            "use_struct_tree": self._config_bool(
                bot_config, "pdf_use_struct_tree", settings.PDF_USE_STRUCT_TREE
            ),
            "include_header_footer": self._config_bool(
                bot_config, "pdf_include_header_footer", settings.PDF_INCLUDE_HEADER_FOOTER
            ),
            "detect_strikethrough": self._config_bool(
                bot_config, "pdf_detect_strikethrough", settings.PDF_DETECT_STRIKETHROUGH
            ),
            "threads": max(1, self._config_int(bot_config, "pdf_threads", settings.PDF_THREADS)),
        }

    def _build_opendataloader_convert_kwargs(
        self,
        file_path: str,
        output_dir: str,
        image_dir: str,
        pdf_config: Dict[str, Any],
    ) -> Dict[str, Any]:
        """Build opendataloader-pdf convert kwargs using the non-deprecated v2.4 API."""
        convert_kwargs: Dict[str, Any] = {
            "input_path": [file_path],
            "output_dir": output_dir,
            "format": "markdown,json",
            "table_method": settings.PDF_TABLE_METHOD,
            "reading_order": settings.PDF_READING_ORDER,
            "image_output": settings.PDF_IMAGE_OUTPUT,
            "image_format": settings.PDF_IMAGE_FORMAT,
            "image_dir": image_dir,
            "markdown_page_separator": settings.PDF_PAGE_SEPARATOR,
            "markdown_with_html": settings.PDF_MARKDOWN_WITH_HTML,
            "sanitize": pdf_config["sanitize"],
            "use_struct_tree": pdf_config["use_struct_tree"],
            "include_header_footer": pdf_config["include_header_footer"],
            "detect_strikethrough": pdf_config["detect_strikethrough"],
            "quiet": True,
        }

        if pdf_config["threads"] > 1:
            convert_kwargs["threads"] = str(pdf_config["threads"])

        parser_mode = pdf_config["parser_mode"]
        hybrid_provider = settings.PDF_HYBRID_MODE
        if parser_mode == "hybrid_full" and hybrid_provider == "off":
            hybrid_provider = "docling-fast"

        if parser_mode != "local_fast" and hybrid_provider != "off":
            convert_kwargs["hybrid"] = hybrid_provider
            convert_kwargs["hybrid_url"] = settings.PDF_HYBRID_URL
            convert_kwargs["hybrid_fallback"] = settings.PDF_HYBRID_FALLBACK
            convert_kwargs["hybrid_mode"] = "full" if parser_mode == "hybrid_full" else "auto"

        return convert_kwargs

    def _load_document(
        self,
        file_path: str,
        filename: str,
        enrich_picture_description: bool = False,
        document_id: Optional[str] = None,
        bot_config: Optional[dict] = None,
    ) -> List[LangChainDocument]:
        """Load document based on file type."""
        suffix = Path(filename).suffix.lower()
        if suffix == ".pdf":
            return self._load_pdf_opendataloader(
                file_path,
                filename,
                enrich_picture_description=enrich_picture_description,
                document_id=document_id,
                bot_config=bot_config,
            )
        elif suffix in {".txt", ".md"}:
            loader = TextLoader(file_path, encoding="utf-8")
            return loader.load()
        elif suffix == ".csv":
            return self._load_csv(file_path, filename)
        elif suffix == ".docx":
            return self._load_docx(file_path, filename)
        elif suffix == ".pptx":
            return self._load_pptx(file_path, filename)
        elif suffix == ".xlsx":
            return self._load_xlsx(file_path, filename)
        elif suffix in {".doc", ".ppt", ".xls"}:
            raise ValueError(
                f"Unsupported legacy Office file type: {suffix}. "
                "Please convert to .docx, .pptx, or .xlsx and upload again."
            )
        else:
            raise ValueError(f"Unsupported file type: {filename}")

    def _load_csv(self, file_path: str, filename: str) -> List[LangChainDocument]:
        rows = []
        with open(file_path, "r", encoding="utf-8-sig", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(cell.strip() for cell in row))
        text = "\n".join(row for row in rows if row.strip())
        return [LangChainDocument(page_content=text, metadata={"source": filename, "page": 0})]

    def _load_docx(self, file_path: str, filename: str) -> List[LangChainDocument]:
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                line = " | ".join(cell for cell in cells if cell)
                if line:
                    parts.append(line)
        return [
            LangChainDocument(
                page_content="\n\n".join(parts),
                metadata={"source": filename, "page": 0},
            )
        ]

    def _load_pptx(self, file_path: str, filename: str) -> List[LangChainDocument]:
        from pptx import Presentation

        prs = Presentation(file_path)
        documents = []
        for idx, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        line = " | ".join(cell for cell in cells if cell)
                        if line:
                            parts.append(line)
            text = "\n\n".join(parts).strip()
            if text:
                documents.append(
                    LangChainDocument(
                        page_content=text,
                        metadata={"source": filename, "page": idx},
                    )
                )
        return documents or [LangChainDocument(page_content="", metadata={"source": filename, "page": 0})]

    def _load_xlsx(self, file_path: str, filename: str) -> List[LangChainDocument]:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        documents = []
        for ws in wb.worksheets:
            lines = []
            for row in ws.iter_rows(values_only=True):
                cells = ["" if value is None else str(value).strip() for value in row]
                line = " | ".join(cell for cell in cells if cell)
                if line:
                    lines.append(line)
            text = "\n".join(lines).strip()
            if text:
                documents.append(
                    LangChainDocument(
                        page_content=f"# Sheet: {ws.title}\n\n{text}",
                        metadata={"source": filename, "sheet": ws.title, "page": 0},
                    )
                )
        wb.close()
        return documents or [LangChainDocument(page_content="", metadata={"source": filename, "page": 0})]

    def _load_pdf_opendataloader(
        self,
        file_path: str,
        filename: str,
        enrich_picture_description: bool = False,
        document_id: Optional[str] = None,
        bot_config: Optional[dict] = None,
    ) -> List[LangChainDocument]:
        """Parse PDF using OpenDataLoader JSON/Markdown artifacts for layout-aware RAG."""
        try:
            import opendataloader_pdf
        except ImportError:
            logger.warning("opendataloader-pdf not installed, falling back to PyPDFLoader")
            from langchain_community.document_loaders import PyPDFLoader
            return PyPDFLoader(file_path).load()

        pdf_config = self._resolve_pdf_runtime_config(
            bot_config=bot_config,
            enrich_picture_description=enrich_picture_description,
        )

        with tempfile.TemporaryDirectory() as output_dir:
            image_dir = os.path.join(output_dir, "images")
            os.makedirs(image_dir, exist_ok=True)

            convert_kwargs = self._build_opendataloader_convert_kwargs(
                file_path=file_path,
                output_dir=output_dir,
                image_dir=image_dir,
                pdf_config=pdf_config,
            )

            logger.info(
                "opendataloader_convert_start",
                extra={
                    "source_filename": filename,
                    "parser_mode": pdf_config["parser_mode"],
                    "hybrid_mode": convert_kwargs.get("hybrid_mode"),
                    "structured_chunking": pdf_config["structured_chunking"],
                },
            )

            try:
                opendataloader_pdf.convert(**convert_kwargs)
            except Exception as convert_err:
                if "hybrid" in convert_kwargs and settings.PDF_HYBRID_FALLBACK:
                    logger.warning(f"Hybrid mode failed ({convert_err}), retrying local mode")
                    convert_kwargs.pop("hybrid", None)
                    convert_kwargs.pop("hybrid_url", None)
                    convert_kwargs.pop("hybrid_fallback", None)
                    convert_kwargs.pop("hybrid_mode", None)
                    try:
                        opendataloader_pdf.convert(**convert_kwargs)
                    except Exception as local_err:
                        logger.warning(f"Local mode also failed ({local_err}), PyPDFLoader fallback")
                        from langchain_community.document_loaders import PyPDFLoader
                        return PyPDFLoader(file_path).load()
                else:
                    raise

            # --- Read Markdown output ---
            stem = Path(file_path).stem
            md_path = os.path.join(output_dir, stem + ".md")
            if not os.path.exists(md_path):
                md_files = list(Path(output_dir).glob("*.md"))
                if not md_files:
                    raise RuntimeError(
                        f"opendataloader-pdf produced no markdown output for {filename}"
                    )
                md_path = str(md_files[0])

            with open(md_path, "r", encoding="utf-8") as f:
                text = f.read()

            if not text.strip():
                logger.warning(
                    f"opendataloader-pdf returned empty content for {filename}, "
                    "falling back to PyPDFLoader"
                )
                from langchain_community.document_loaders import PyPDFLoader
                return PyPDFLoader(file_path).load()

            logger.info(f"opendataloader-pdf parsed {filename}: {len(text)} chars")

            json_path = os.path.join(output_dir, stem + ".json")
            if not os.path.exists(json_path):
                json_files = list(Path(output_dir).glob("*.json"))
                json_path = str(json_files[0]) if json_files else ""

            artifact_paths = self._upload_opendataloader_artifacts(
                md_path=md_path,
                json_path=json_path if json_path and os.path.exists(json_path) else None,
                image_dir=image_dir,
                original_filename=filename,
                document_id=document_id,
            )

            if artifact_paths.get("markdown"):
                try:
                    with open(md_path, "r", encoding="utf-8") as f:
                        text = f.read()
                except OSError:
                    pass

            if pdf_config["structured_chunking"] and json_path and os.path.exists(json_path):
                try:
                    with open(json_path, "r", encoding="utf-8") as f:
                        json_data = json.load(f)
                    structured_docs = self._documents_from_opendataloader_json(
                        json_data=json_data,
                        source=filename,
                        artifact_paths=artifact_paths,
                    )
                    if structured_docs:
                        logger.info(
                            "opendataloader_structured_docs_created",
                            extra={"source_filename": filename, "documents": len(structured_docs)},
                        )
                        return structured_docs
                except Exception as e:
                    logger.warning(f"OpenDataLoader JSON normalization failed for {filename}: {e}")

            doc_metadata = {
                "source": filename,
                "page": 0,
                "page_numbers": [],
                "bboxes": [],
                "element_types": ["markdown"],
                "heading_path": [],
                "opendataloader_element_ids": [],
                "has_structured_json": bool(json_path and os.path.exists(json_path)),
                "artifact_paths": artifact_paths,
            }

            return [
                LangChainDocument(
                    page_content=text,
                    metadata=doc_metadata,
                )
            ]

    def _upload_opendataloader_artifacts(
        self,
        md_path: str,
        json_path: Optional[str],
        image_dir: str,
        original_filename: str,
        document_id: Optional[str],
    ) -> Dict[str, Any]:
        """Upload extracted Markdown, JSON, and images under a stable MinIO prefix."""
        artifact_paths: Dict[str, Any] = {
            "prefix": f"documents/{document_id}/extracted" if document_id else f"documents/{uuid.uuid4()}/extracted",
            "images": [],
        }
        try:
            from app.services.storage_service import storage_service
            prefix = artifact_paths["prefix"]

            image_files = [
                p for p in Path(image_dir).glob("*.*")
                if p.suffix.lower() in (".png", ".jpeg", ".jpg")
            ]
            link_rewrites: Dict[str, str] = {}
            for img_path in image_files:
                content_type = "image/png" if img_path.suffix.lower() == ".png" else "image/jpeg"
                object_path = f"{prefix}/images/{img_path.name}"
                with open(img_path, "rb") as f:
                    storage_service.upload_file_to_path(f, object_path, content_type)
                artifact_paths["images"].append(object_path)
                rel_path = f"{Path(image_dir).name}/{img_path.name}"
                link_rewrites[rel_path] = f"minio://{settings.MINIO_BUCKET}/{object_path}"
                link_rewrites[f"./{rel_path}"] = f"minio://{settings.MINIO_BUCKET}/{object_path}"
                link_rewrites[img_path.name] = f"minio://{settings.MINIO_BUCKET}/{object_path}"

            if link_rewrites:
                with open(md_path, "r", encoding="utf-8") as f:
                    md_text = f.read()
                for rel_path, target in link_rewrites.items():
                    md_text = md_text.replace(f"]({rel_path})", f"]({target})")
                    md_text = md_text.replace(f"](<{rel_path}>)", f"](<{target}>)")
                with open(md_path, "w", encoding="utf-8") as f:
                    f.write(md_text)

            md_filename = Path(original_filename).stem + ".md"
            with open(md_path, "rb") as f:
                artifact_paths["markdown"] = storage_service.upload_file_to_path(
                    f, f"{prefix}/{md_filename}", content_type="text/markdown"
                )

            if json_path:
                json_filename = Path(original_filename).stem + ".json"
                with open(json_path, "rb") as f:
                    artifact_paths["json"] = storage_service.upload_file_to_path(
                        f, f"{prefix}/{json_filename}", content_type="application/json"
                    )

            logger.info(
                "opendataloader_artifacts_uploaded",
                extra={"source_filename": original_filename, "prefix": prefix, "images": len(artifact_paths["images"])},
            )
        except Exception as e:
            logger.warning(f"Failed to upload OpenDataLoader artifacts to MinIO: {e}")
        return artifact_paths

    @staticmethod
    def _coerce_bbox(value: Any) -> Optional[List[float]]:
        if value is None:
            return None
        if isinstance(value, (list, tuple)) and len(value) >= 4:
            try:
                return [float(value[i]) for i in range(4)]
            except (TypeError, ValueError):
                return None
        if isinstance(value, dict):
            candidates = [
                ("left", "bottom", "right", "top"),
                ("l", "b", "r", "t"),
                ("x0", "y0", "x1", "y1"),
            ]
            for keys in candidates:
                if all(k in value for k in keys):
                    try:
                        return [float(value[k]) for k in keys]
                    except (TypeError, ValueError):
                        return None
        return None

    @classmethod
    def _extract_element_page_bbox(cls, element: Dict[str, Any]) -> tuple[Optional[int], Optional[List[float]]]:
        page = (
            element.get("page number")
            or element.get("page_number")
            or element.get("page")
            or element.get("page_no")
        )
        bbox = element.get("bounding box") or element.get("bounding_box") or element.get("bbox")
        prov = element.get("prov")
        if isinstance(prov, list) and prov:
            first = prov[0] if isinstance(prov[0], dict) else {}
            page = page or first.get("page_no") or first.get("page")
            bbox = bbox or first.get("bbox")

        try:
            page_int = int(page) if page is not None else None
        except (TypeError, ValueError):
            page_int = None
        return page_int, cls._coerce_bbox(bbox)

    @staticmethod
    def _table_data_to_text(table_data: Any) -> str:
        if not table_data:
            return ""
        if isinstance(table_data, str):
            return table_data
        if not isinstance(table_data, dict):
            return ""

        grid = table_data.get("grid") or table_data.get("rows")
        if isinstance(grid, list):
            rows: List[str] = []
            for row in grid:
                cells = []
                if isinstance(row, list):
                    for cell in row:
                        if isinstance(cell, dict):
                            cells.append(str(cell.get("text") or cell.get("content") or "").strip())
                        else:
                            cells.append(str(cell).strip())
                elif isinstance(row, dict):
                    cells.append(str(row.get("text") or row.get("content") or "").strip())
                line = " | ".join(cell for cell in cells if cell)
                if line:
                    rows.append(line)
            if rows:
                return "\n".join(rows)

        table_cells = table_data.get("table_cells") or table_data.get("cells")
        if isinstance(table_cells, list):
            rows_by_idx: Dict[int, List[tuple[int, str]]] = {}
            for cell in table_cells:
                if not isinstance(cell, dict):
                    continue
                row_idx = int(cell.get("start_row_offset_idx") or cell.get("row") or 0)
                col_idx = int(cell.get("start_col_offset_idx") or cell.get("col") or 0)
                text = str(cell.get("text") or cell.get("content") or "").strip()
                if text:
                    rows_by_idx.setdefault(row_idx, []).append((col_idx, text))
            rows = []
            for row_idx in sorted(rows_by_idx):
                rows.append(" | ".join(text for _, text in sorted(rows_by_idx[row_idx])))
            return "\n".join(rows)
        return ""

    @classmethod
    def _extract_element_text(cls, element: Dict[str, Any], element_type: str) -> str:
        content = (
            element.get("content")
            or element.get("text")
            or element.get("orig")
            or element.get("alt")
            or element.get("description")
            or ""
        )
        if not content and element_type == "table":
            content = cls._table_data_to_text(element.get("data") or element)
        if not content and element_type in {"image", "picture"}:
            annotations = element.get("annotations") or element.get("captions") or []
            if isinstance(annotations, list):
                parts = []
                for item in annotations:
                    if isinstance(item, dict):
                        parts.append(str(item.get("text") or item.get("content") or "").strip())
                    elif isinstance(item, str):
                        parts.append(item.strip())
                content = " ".join(part for part in parts if part)
        content = str(content).strip()
        if not content:
            return ""
        if element_type == "table":
            return f"Table:\n{content}"
        if element_type == "formula":
            return f"Formula: {content}"
        if element_type in {"image", "picture"}:
            return f"Image description: {content}"
        return content

    @classmethod
    def _documents_from_opendataloader_json(
        cls,
        json_data: Any,
        source: str,
        artifact_paths: Optional[Dict[str, Any]] = None,
    ) -> List[LangChainDocument]:
        """Normalize OpenDataLoader JSON into element-level LangChain documents."""
        documents: List[LangChainDocument] = []
        heading_path: List[str] = []

        def visit(element: Any, inherited_page: Optional[int] = None):
            nonlocal heading_path
            if not isinstance(element, dict):
                return

            element_type = str(
                element.get("type") or element.get("label") or element.get("kind") or "text"
            ).replace("-", "_").lower()
            text = cls._extract_element_text(element, element_type)
            page, bbox = cls._extract_element_page_bbox(element)
            page = page or inherited_page

            if element_type in {"heading", "section_header", "title"} and text:
                try:
                    level = int(element.get("heading level") or element.get("heading_level") or element.get("level") or 1)
                except (TypeError, ValueError):
                    level = 1
                level = max(1, min(level, 10))
                heading_path = heading_path[: level - 1] + [text]

            if text and element_type not in {"document"}:
                element_id = element.get("id") or element.get("self_ref") or element.get("ref")
                metadata = {
                    "source": source,
                    "page": page or 0,
                    "page_numbers": [page] if page else [],
                    "bboxes": [bbox] if bbox else [],
                    "element_types": [element_type],
                    "heading_path": heading_path.copy(),
                    "opendataloader_element_ids": [str(element_id)] if element_id else [],
                    "has_structured_json": True,
                    "artifact_paths": artifact_paths or {},
                }
                documents.append(LangChainDocument(page_content=text, metadata=metadata))

            for child_key in ("kids", "elements"):
                children = element.get(child_key)
                if isinstance(children, list):
                    for child in children:
                        visit(child, page)

        if isinstance(json_data, dict):
            if isinstance(json_data.get("kids"), list):
                for item in json_data["kids"]:
                    visit(item)
            pages = json_data.get("pages")
            if isinstance(pages, dict):
                for page_key in sorted(pages, key=lambda value: int(value) if str(value).isdigit() else str(value)):
                    page_data = pages[page_key]
                    try:
                        page_num = int(page_key)
                    except (TypeError, ValueError):
                        page_num = None
                    if isinstance(page_data, dict):
                        for key in ("elements", "kids"):
                            if isinstance(page_data.get(key), list):
                                for item in page_data[key]:
                                    visit(item, page_num)
            elif isinstance(pages, list):
                for idx, page_data in enumerate(pages, start=1):
                    if isinstance(page_data, dict):
                        page_num = page_data.get("page") or page_data.get("page_number") or idx
                        for key in ("elements", "kids"):
                            if isinstance(page_data.get(key), list):
                                for item in page_data[key]:
                                    visit(item, int(page_num) if page_num else idx)
            if not documents:
                visit(json_data)
        elif isinstance(json_data, list):
            for item in json_data:
                visit(item)

        return [doc for doc in documents if doc.page_content.strip()]
    
    @staticmethod
    def _chunk_article(text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
        """
        Article-boundary chunking for legal documents.
        Splits at Vietnamese article markers (Điều N) then sub-splits oversized articles.
        """
        # Split at article boundaries — keep the "Điều" marker with the following text
        parts = re.split(r'(?=Điều\s+\d+)', text)
        result: List[str] = []
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        for part in parts:
            part = part.strip()
            if not part:
                continue
            if len(part) <= chunk_size:
                result.append(part)
            else:
                # Sub-split oversized articles recursively
                result.extend(splitter.split_text(part))
        return result

    @staticmethod
    def _chunk_sentence(text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
        """
        Sentence-aware chunking using a rolling-window accumulator.
        Preserves paragraph coherence by joining sentences up to chunk_size.
        """
        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
        chunks: List[str] = []
        current: List[str] = []
        current_len = 0
        for sent in sentences:
            sent_len = len(sent)
            if current_len + sent_len > chunk_size and current:
                chunks.append(" ".join(current))
                # Keep overlap by retaining last few sentences
                overlap_chars = 0
                keep: List[str] = []
                for s in reversed(current):
                    overlap_chars += len(s)
                    keep.insert(0, s)
                    if overlap_chars >= chunk_overlap:
                        break
                current = keep
                current_len = sum(len(s) for s in current)
            current.append(sent)
            current_len += sent_len
        if current:
            chunks.append(" ".join(current))
        return chunks

    def _chunk_documents(
        self,
        documents: List[LangChainDocument],
        chunking_strategy: str = "recursive",
        chunk_size: int = 1000,
        chunk_overlap: int = 200
    ) -> List[LangChainDocument]:
        """Split documents into chunks using specified strategy."""

        if chunking_strategy == "article":
            chunks: List[LangChainDocument] = []
            for doc in documents:
                for piece in self._chunk_article(doc.page_content, chunk_size, chunk_overlap):
                    chunks.append(LangChainDocument(page_content=piece, metadata=doc.metadata.copy()))
            logger.info(f"Created {len(chunks)} chunks using article strategy")
            return chunks

        if chunking_strategy == "sentence":
            chunks = []
            for doc in documents:
                for piece in self._chunk_sentence(doc.page_content, chunk_size, chunk_overlap):
                    chunks.append(LangChainDocument(page_content=piece, metadata=doc.metadata.copy()))
            logger.info(f"Created {len(chunks)} chunks using sentence strategy")
            return chunks

        if chunking_strategy == "parent_child":
            # Parent-Child Chunking:
            # Large parent chunks preserve context; small child chunks enable precise matching.
            # Each child carries its parent text in metadata so retrieval returns the richer context.
            parent_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                separators=["\n\n", "\n", ". ", " ", ""],
            )
            child_splitter = RecursiveCharacterTextSplitter(
                chunk_size=max(128, chunk_size // 4),
                chunk_overlap=max(32, chunk_overlap // 2),
                separators=["\n\n", "\n", ". ", " ", ""],
            )
            result_chunks: List[LangChainDocument] = []
            for doc in documents:
                for parent_text in parent_splitter.split_text(doc.page_content):
                    for child_piece in child_splitter.split_text(parent_text):
                        meta = doc.metadata.copy()
                        meta["parent_text"] = parent_text
                        result_chunks.append(LangChainDocument(page_content=child_piece, metadata=meta))
            logger.info(f"Created {len(result_chunks)} child chunks with parent_child strategy")
            return result_chunks

        if chunking_strategy == "semantic":
            # Semantic chunking: split at sentence boundaries
            try:
                from langchain_text_splitters import SpacyTextSplitter
                text_splitter = SpacyTextSplitter(
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                    pipeline="en_core_web_sm"
                )
                logger.info("Using semantic chunking with spacy")
            except Exception as e:
                logger.warning(f"Spacy not available ({e}), falling back to recursive")
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                    separators=["\n\n", "\n", ". ", " ", ""]
                )
        else:  # recursive (default)
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                separators=["\n\n", "\n", ". ", " ", ""]
            )

        chunks = text_splitter.split_documents(documents)
        logger.info(f"Created {len(chunks)} chunks using {chunking_strategy} strategy")
        return chunks
    
    @retry(
        retry=retry_if_exception_type((httpx.HTTPError, OpenRouterAPIError, Exception)),
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=2, max=10),
        reraise=True
    )
    def _embed_with_retry(self, text: str) -> List[float]:
        """Embed single text with retry logic"""
        try:
            return self.openrouter.embed_single(text)
        except Exception as e:
            logger.warning(f"Embedding attempt failed: {e}")
            raise OpenRouterAPIError(f"Embedding failed: {str(e)}")
    
    @retry(
        retry=retry_if_exception_type((httpx.HTTPError, OpenRouterAPIError, Exception)),
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=2, max=10),
        reraise=True
    )
    def _chat_with_retry(self, messages: List[Dict], **kwargs) -> Dict:
        """Chat completion with retry logic"""
        try:
            return self.openrouter.chat_completion(messages=messages, **kwargs)
        except Exception as e:
            logger.warning(f"Chat completion attempt failed: {e}")
            raise OpenRouterAPIError(f"Chat completion failed: {str(e)}")
    
    def _get_reranker(self):
        """Lazy loader for the reranker model to speed up startup.

        Device priority: CUDA → MPS (Apple Silicon, native macOS only) → CPU.
        Note: MPS is unavailable inside Docker on Mac (Linux VM has no Metal access).
        To use MPS/GPU on M1/M2/M3, run backend natively: uvicorn app.main:app --port 8000
        """
        if self._reranker_attempted:
            return self.reranker

        self._reranker_attempted = True
        try:
            import torch
            from sentence_transformers import CrossEncoder

            model_cache_dir = os.getenv('HF_HOME', '/tmp/huggingface_cache')
            os.makedirs(model_cache_dir, exist_ok=True)

            reranker_model = getattr(settings, "RERANKER_MODEL", "cross-encoder/ms-marco-MiniLM-L-6-v2")

            # Auto-detect best available device
            if torch.cuda.is_available():
                device = "cuda"
            elif torch.backends.mps.is_available():
                device = "mps"
            else:
                device = "cpu"

            logger.info("reranker_loading", extra={"model": reranker_model, "device": device.upper()})
            self.reranker = CrossEncoder(
                reranker_model,
                device=device,
            )
            logger.info("reranker_loaded", extra={"device": device.upper()})
        except Exception as e:
            logger.warning(f"Failed to load Reranker model on demand: {e}")
            self.reranker = None

        return self.reranker

    async def _hybrid_search(
        self,
        bot_id: str,
        query: str,
        query_embedding: List[float],
        top_k: int = 5,
        rerank: bool = True,
    ) -> List[Dict]:
        """
        True hybrid retrieval: dense OpenRouter vectors + sparse BM25 vectors
        fused inside Qdrant via RRF, then optionally re-ranked with Cross-Encoder.

        rerank=False skips the Cross-Encoder pass (used for multi-query variants where
        the caller does one final rerank after merging all variant results).
        """
        reranker = self._get_reranker() if rerank else None
        initial_limit = max(top_k * 4, 20)
        
        bot_filter = Filter(
            must=[FieldCondition(key="bot_id", match=MatchValue(value=bot_id))]
        )

        try:
            sparse_query = await asyncio.to_thread(self._embed_sparse_query, query)
            response = await self.async_qdrant_client.query_points(
                collection_name=self.collection_name,
                prefetch=[
                    rest.Prefetch(
                        query=query_embedding,
                        using=self.dense_vector_name,
                        limit=initial_limit,
                    ),
                    rest.Prefetch(
                        query=sparse_query,
                        using=self.sparse_vector_name,
                        limit=initial_limit,
                    ),
                ],
                query=self._build_rrf_query(),
                query_filter=bot_filter,
                limit=initial_limit,
                with_payload=True,
                with_vectors=False,
            )
            points = response.points
            logger.debug(f"Qdrant dense+sparse RRF found {len(points)} candidates for query: {query[:50]}")
        except Exception as e:
            logger.warning(f"Qdrant hybrid search failed, falling back to dense-only: {e}")
            response = await self.async_qdrant_client.query_points(
                collection_name=self.collection_name,
                query=query_embedding,
                using=self.dense_vector_name,
                query_filter=bot_filter,
                limit=initial_limit,
                with_payload=True,
                with_vectors=False,
            )
            points = response.points

        candidates = []
        for point in points:
            payload = point.payload or {}
            candidates.append({
                "id": str(point.id),
                "text": payload.get("text", ""),
                "parent_text": payload.get("parent_text"),
                "source": payload.get("source", "unknown"),
                "document_id": payload.get("document_id"),
                "context_prefix": payload.get("context_prefix"),
                "initial_score": point.score,
                "rrf_score": point.score,
                "metadata": payload.get("metadata", {}),
            })

        if not candidates:
            return []

        # 4. Re-ranking (Cross-Encoder) if available
        import numpy as np
        if reranker:
            try:
                pairs = [[query, c["text"]] for c in candidates]
                rerank_scores = reranker.predict(pairs)

                if isinstance(rerank_scores, float):
                    rerank_scores = [rerank_scores]

                norm_scores = _sigmoid(np.array(rerank_scores))

                for i, c in enumerate(candidates):
                    c["hybrid_score"] = float(norm_scores[i])
                    c["rerank_raw"] = float(rerank_scores[i])

                candidates.sort(key=lambda x: x["hybrid_score"], reverse=True)
                logger.info(f"Hybrid reranked {len(candidates)} items (dense+bm25+ce). Top: {candidates[0]['hybrid_score']:.4f}")

            except Exception as e:
                logger.warning(f"Reranking failed: {e}. Falling back to RRF order.")
                for c in candidates:
                    c["hybrid_score"] = c["rrf_score"]
        else:
            for c in candidates:
                c["hybrid_score"] = c["rrf_score"]
            candidates.sort(key=lambda x: x["hybrid_score"], reverse=True)

        return candidates[:top_k]

    async def _generate_hyde_query(self, query: str) -> str:
        """
        Generate hypothetical document using HyDE technique.
        Creates a synthetic answer that's embedded for better retrieval.
        """
        hyde_prompt = [
            {
                "role": "system",
                "content": "You are a helpful assistant that generates detailed, factual answers."
            },
            {
                "role": "user",
                "content": f"""Write a detailed, comprehensive answer to the following question. 
Even if you don't know the exact answer, write what a good answer would look like based on the question.

Question: {query}

Answer:"""
            }
        ]

        try:
            response = self._chat_with_retry(
                messages=hyde_prompt,
                temperature=0.3,
                max_tokens=200
            )
            hyde_doc = response["content"]
            logger.info(f"Generated HyDE document: {hyde_doc[:100]}...")
            return hyde_doc
        except Exception as e:
            logger.warning(f"HyDE generation failed, using original query: {e}")
            return query
    
    async def ingest_file(
        self,
        file: UploadFile,
        bot_id: str,
        chunking_strategy: str = "recursive",
        chunk_size: int = 1000,
        chunk_overlap: int = 200
    ) -> Dict[str, Any]:
        """
        Ingest file into vector database using OpenRouter embeddings.
        
        Args:
            file: Uploaded file
            bot_id: Bot ID for filtering
            chunking_strategy: Chunking strategy to use
            chunk_size: Size of each chunk
            chunk_overlap: Overlap between chunks
            
        Returns:
            Dict with ingestion stats
        """
        # Save temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{file.filename}") as tmp:
            shutil.copyfileobj(file.file, tmp)
            tmp_path = tmp.name
        
        try:
            start_time = time.time()
            
            # Load document
            logger.info(f"Loading document: {file.filename}")
            documents = self._load_document(tmp_path, file.filename)
            
            # Chunk documents
            logger.info(f"Chunking document with strategy={chunking_strategy}, size={chunk_size}, overlap={chunk_overlap}")
            chunks = self._chunk_documents(
                documents, 
                chunking_strategy=chunking_strategy,
                chunk_size=chunk_size, 
                chunk_overlap=chunk_overlap
            )
            
            # Add metadata
            for chunk in chunks:
                chunk.metadata["bot_id"] = bot_id
                chunk.metadata["source"] = file.filename
                chunk.metadata["ingested_at"] = datetime.utcnow().isoformat()

            # ── Contextual Retrieval: generate situating prefix per chunk ──────
            chunk_texts = [chunk.page_content for chunk in chunks]
            doc_text_full = "\n\n".join(doc.page_content for doc in documents)
            logger.info(f"Generating contextual prefixes for {len(chunks)} chunks (Contextual Retrieval)")
            context_prefixes = await self._generate_contextual_prefix_batch(doc_text_full, chunk_texts)
            enriched_texts = [
                f"{prefix}\n\n{text}" if prefix else text
                for prefix, text in zip(context_prefixes, chunk_texts)
            ]
            # ─────────────────────────────────────────────────────────────────

            # Generate dense + sparse embeddings for all chunks.
            logger.info(f"Generating dense+sparse embeddings for {len(chunks)} chunks")
            embeddings, sparse_embeddings = await asyncio.gather(
                self.openrouter.embed_batch_async(enriched_texts, batch_size=100),
                asyncio.to_thread(self._embed_sparse_texts, enriched_texts),
            )
            if len(embeddings) != len(chunks) or len(sparse_embeddings) != len(chunks):
                raise OpenRouterAPIError(
                    f"Embedding count mismatch: chunks={len(chunks)}, dense={len(embeddings)}, sparse={len(sparse_embeddings)}"
                )
            
            # Auto-detect embedding dimensions
            if embeddings:
                detected_dim = len(embeddings[0])
                if detected_dim != self.embedding_dim:
                    raise OpenRouterAPIError(
                        f"Embedding dimension mismatch. Expected {self.embedding_dim}, "
                        f"got {detected_dim}. Refusing to mutate Qdrant collection."
                    )
            
            # Insert into Qdrant
            logger.info(f"Inserting {len(chunks)} vectors into Qdrant")
            points = []
            for idx, (chunk, embedding, sparse_embedding) in enumerate(zip(chunks, embeddings, sparse_embeddings)):
                point_id = hashlib.md5(
                    f"{bot_id}_{file.filename}_{idx}_{chunk.page_content[:100]}".encode()
                ).hexdigest()

                points.append(
                    PointStruct(
                        id=point_id,
                        vector=self._build_named_vector(embedding, sparse_embedding),
                        payload={
                            "bot_id": bot_id,
                            "document_id": None,
                            "source": file.filename,
                            "text": chunk.page_content,
                            "parent_text": chunk.metadata.get("parent_text"),
                            "context_prefix": context_prefixes[idx] if idx < len(context_prefixes) else None,
                            "page_numbers": chunk.metadata.get("page_numbers", []),
                            "bboxes": chunk.metadata.get("bboxes", []),
                            "element_types": chunk.metadata.get("element_types", []),
                            "artifact_paths": chunk.metadata.get("artifact_paths", {}),
                            "metadata": chunk.metadata,
                        }
                    )
                )
            
            await self.async_qdrant_client.upsert(
                collection_name=self.collection_name,
                points=points,
            )
            
            elapsed_time = time.time() - start_time
            
            result = {
                "filename": file.filename,
                "chunks_created": len(chunks),
                "vectors_inserted": len(points),
                "embedding_dim": self.embedding_dim,
                "processing_time": round(elapsed_time, 2)
            }
            
            logger.info(f"File ingestion complete: {result}")
            
            # Invalidate cache for this bot since knowledge base changed
            await self.invalidate_bot_cache_async(bot_id)
            
            return result
            
        finally:
            # Clean up temp file
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    def process_file_sync(
        self,
        file_path: str,
        bot_id: str,
        filename: str,
        chunking_strategy: str = "recursive",
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        preloaded_documents=None,
        document_id: Optional[str] = None,
        bot_config: Optional[dict] = None,
    ) -> Dict[str, Any]:
        """
        Synchronous file ingestion for Celery worker.

        Pass ``preloaded_documents`` to skip the internal ``_load_document``
        call when the caller has already parsed the file (e.g. to also feed
        the same content to LightRAG), avoiding a redundant disk read.
        """
        start_time = time.time()

        if preloaded_documents is not None:
            documents = preloaded_documents
            logger.info(f"Using pre-loaded document for: {filename}")
        else:
            logger.info(f"Loading document: {filename}")
            documents = self._load_document(
                file_path,
                filename,
                document_id=document_id,
                bot_config=bot_config,
            )

        logger.info(
            f"Chunking document with strategy={chunking_strategy}, size={chunk_size}, overlap={chunk_overlap}"
        )
        chunks = self._chunk_documents(
            documents,
            chunking_strategy=chunking_strategy,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )

        for chunk in chunks:
            chunk.metadata["bot_id"] = bot_id
            chunk.metadata["source"] = filename
            chunk.metadata["ingested_at"] = datetime.utcnow().isoformat()

        # ── Contextual Retrieval + embedding (single async pass) ──────────────
        async def _contextual_ingest_async() -> tuple:
            doc_text_full = "\n\n".join(doc.page_content for doc in documents)
            chunk_texts = [c.page_content for c in chunks]
            logger.info(f"Generating contextual prefixes for {len(chunks)} chunks (Contextual Retrieval)")
            prefixes = await self._generate_contextual_prefix_batch(doc_text_full, chunk_texts)
            enriched = [
                f"{p}\n\n{t}" if p else t for p, t in zip(prefixes, chunk_texts)
            ]
            dense_task = self.openrouter.embed_batch_async(enriched, batch_size=100)
            sparse_task = asyncio.to_thread(self._embed_sparse_texts, enriched)
            embeddings, sparse_embeddings = await asyncio.gather(dense_task, sparse_task)
            return prefixes, embeddings, sparse_embeddings

        logger.info(f"Generating dense+sparse embeddings for {len(chunks)} chunks")
        context_prefixes, embeddings, sparse_embeddings = asyncio.run(_contextual_ingest_async())
        if len(embeddings) != len(chunks) or len(sparse_embeddings) != len(chunks):
            raise OpenRouterAPIError(
                f"Embedding count mismatch: chunks={len(chunks)}, dense={len(embeddings)}, sparse={len(sparse_embeddings)}"
            )
        # ─────────────────────────────────────────────────────────────────────

        if embeddings:
            detected_dim = len(embeddings[0])
            if detected_dim != self.embedding_dim:
                raise OpenRouterAPIError(
                    f"Embedding dimension mismatch. Expected {self.embedding_dim}, "
                    f"got {detected_dim}. Refusing to mutate Qdrant collection."
                )

        points = []
        for idx, (chunk, embedding, sparse_embedding) in enumerate(zip(chunks, embeddings, sparse_embeddings)):
            point_id = hashlib.md5(
                f"{bot_id}_{document_id or filename}_{idx}_{chunk.page_content[:100]}".encode()
            ).hexdigest()

            points.append(
                PointStruct(
                    id=point_id,
                    vector=self._build_named_vector(embedding, sparse_embedding),
                    payload={
                        "bot_id": bot_id,
                        "document_id": document_id,
                        "source": filename,
                        "text": chunk.page_content,
                        "parent_text": chunk.metadata.get("parent_text"),
                        "context_prefix": context_prefixes[idx] if idx < len(context_prefixes) else None,
                        "page_numbers": chunk.metadata.get("page_numbers", []),
                        "bboxes": chunk.metadata.get("bboxes", []),
                        "element_types": chunk.metadata.get("element_types", []),
                        "artifact_paths": chunk.metadata.get("artifact_paths", {}),
                        "metadata": chunk.metadata,
                    }
                )
            )

        UPSERT_BATCH_SIZE = 200
        logger.info(f"Inserting {len(points)} vectors into Qdrant (batches of {UPSERT_BATCH_SIZE})")
        batches = [points[i:i + UPSERT_BATCH_SIZE] for i in range(0, len(points), UPSERT_BATCH_SIZE)]
        from app.services.qdrant_gateway import get_qdrant_gateway
        gw = get_qdrant_gateway()

        async def _upsert_batches():
            await asyncio.gather(*[
                gw.upsert(self.collection_name, batch)
                for batch in batches
            ])

        asyncio.run(_upsert_batches())

        elapsed_time = time.time() - start_time

        result = {
            "filename": filename,
            "chunks_created": len(chunks),
            "vectors_inserted": len(points),
            "embedding_dim": self.embedding_dim,
            "processing_time": round(elapsed_time, 2),
            "preview": chunks[0].page_content[:500] if chunks else "No content extracted",
            "model_used": self.openrouter.embedding_model
        }

        logger.info(f"File ingestion complete: {result}")
        self.invalidate_bot_cache(bot_id)

        return result
    
    async def _rewrite_query(self, query: str) -> str:
        """
        Agentic Layer: Rewrite user query to be search-engine friendly.
        Based on the notebook logic.
        """
        system_prompt = """
        Bạn là một chuyên gia tối ưu hóa tìm kiếm (SEO) cho hệ thống RAG.
        Nhiệm vụ của bạn là viết lại câu hỏi của người dùng thành một câu truy vấn tìm kiếm ngắn gọn, súc tích và chứa nhiều từ khóa quan trọng.
        Giữ nguyên ý nghĩa, chỉ thay đổi cấu trúc để dễ tìm kiếm hơn.
        KHÔNG TRẢ LỜI CÂU HỎI. CHỈ VIẾT LẠI QUERY.
        
        Ví dụ:
        - User: "cái app này đăng nhập không được, nó báo lỗi tùm lum tà la"
        - Rewrite: "lỗi không đăng nhập được ứng dụng mobile báo lỗi hệ thống"
        - User: "figure 2 là cái gì"
        - Rewrite: "chi tiết mô tả nội dung Figure 2 hình ảnh 2 trong tài liệu"
        """
        
        try:
            # Quick call to LLM
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": query}
            ]
            
            # Use asyncio.to_thread for synchronous generic call
            response = await asyncio.to_thread(
                self.openrouter.chat_completion,
                messages=messages,
                model=INTERNAL_LLM_MODEL,
                temperature=0.1
            )
            
            rewritten_query = response.get("content", "").strip()
            logger.info(f"Query Rewriting: '{query}' -> '{rewritten_query}'")
            return rewritten_query
            
        except Exception as e:
            logger.warning(f"Query rewriting failed: {e}. Using original query.")
            return query

    async def _generate_hyde_hypothesis(self, query: str, domain: str = "general") -> str:
        """
        HyDE (Hypothetical Document Embedding):
        Generate a hypothetical document passage that *would* answer the query,
        then embed that passage instead of the raw query.

        Doc-to-doc similarity >> query-to-doc similarity in embedding space,
        so this significantly improves semantic retrieval precision.
        Falls back to original query if LLM call fails.
        """
        domain_hints = {
            "legal": (
                "Viết một đoạn văn pháp lý ngắn (4-6 câu) như thể trích từ một văn bản "
                "quy phạm pháp luật, hợp đồng hoặc hướng dẫn pháp lý có thể trả lời câu hỏi sau. "
                "Dùng ngôn ngữ pháp lý chính thức, trích dẫn điều khoản nếu phù hợp."
            ),
            "education": (
                "Viết một đoạn giải thích học thuật ngắn (4-6 câu) như thể trích từ sách giáo khoa "
                "hoặc tài liệu giảng dạy có thể trả lời câu hỏi sau. "
                "Dùng ngôn ngữ rõ ràng, có ví dụ minh họa nếu phù hợp."
            ),
            "sales": (
                "Viết một đoạn mô tả sản phẩm/dịch vụ ngắn (4-6 câu) như thể trích từ tài liệu "
                "kinh doanh, brochure hoặc catalog có thể trả lời câu hỏi sau. "
                "Nhấn mạnh vào lợi ích và đặc điểm nổi bật."
            ),
            "general": (
                "Viết một đoạn văn ngắn (4-6 câu) như thể trích từ tài liệu tham khảo "
                "có thể trả lời câu hỏi sau."
            ),
        }

        hint = domain_hints.get(domain, domain_hints["general"])
        messages = [
            {
                "role": "system",
                "content": (
                    "Bạn là trợ lý tạo đoạn văn giả định để cải thiện tìm kiếm ngữ nghĩa (HyDE). "
                    "Chỉ trả về đoạn văn thuần túy, không thêm giải thích, không thêm tiêu đề."
                ),
            },
            {
                "role": "user",
                "content": f"{hint}\n\nCâu hỏi: {query}",
            },
        ]

        try:
            response = await asyncio.to_thread(
                self.openrouter.chat_completion,
                messages=messages,
                model=INTERNAL_LLM_MODEL,
                temperature=0.5,
                max_tokens=250,
            )
            hypothesis = response.get("content", "").strip()
            if hypothesis:
                logger.info(f"HyDE hypothesis ({domain}): {hypothesis[:120]}...")
                return hypothesis
        except Exception as e:
            logger.warning(f"HyDE generation failed, falling back to original query: {e}")

        return query

    async def _crag_classify(self, query: str, chunks: list) -> str:
        """
        CRAG (Corrective RAG) relevance classifier.

        Takes the top retrieved chunks and the search query, asks a fast LLM to
        decide whether the knowledge base actually contains an answer.

        Returns one of:
          "relevant"   — KB has good information, proceed normally
          "ambiguous"  — partial match, answer with caution
          "no_context" — KB lacks relevant info; LLM should admit it doesn't know

        Falls back to "relevant" on any failure so CRAG never silently breaks RAG.
        """
        if not chunks:
            return "no_context"

        # Build a compact snippet from the top 3 chunks (≤ 300 chars each)
        snippet_lines = []
        for i, chunk in enumerate(chunks[:3], 1):
            text = chunk.get("text", "")[:300].replace("\n", " ")
            snippet_lines.append(f"[{i}] {text}")
        snippets = "\n".join(snippet_lines)

        messages = [
            {
                "role": "system",
                "content": (
                    "You are a relevance judge for a RAG system. "
                    "Given a user question and retrieved document snippets, "
                    "classify the retrieval quality with a single word:\n"
                    "- 'relevant'   → snippets directly answer the question\n"
                    "- 'ambiguous'  → snippets are partially related but incomplete\n"
                    "- 'no_context' → snippets are unrelated or the KB lacks this information\n"
                    "Reply with ONLY one of those three words, nothing else."
                ),
            },
            {
                "role": "user",
                "content": f"Question: {query}\n\nRetrieved snippets:\n{snippets}",
            },
        ]

        try:
            response = await asyncio.to_thread(
                self.openrouter.chat_completion,
                messages=messages,
                model=INTERNAL_LLM_MODEL,
                temperature=0.0,
                max_tokens=16,
            )
            verdict = response.get("content", "").strip().lower()
            if verdict in ("relevant", "ambiguous", "no_context"):
                logger.info(f"CRAG verdict: {verdict}")
                return verdict
        except Exception as e:
            logger.warning(f"CRAG classification failed, defaulting to 'relevant': {e}")

        return "relevant"

    # ──────────────────────────────────────────────────────────────────────────
    # Multi-Query Fusion
    # ──────────────────────────────────────────────────────────────────────────

    async def _generate_query_variants(self, query: str, n: int = 2) -> List[str]:
        """
        Multi-Query Fusion — step 1:
        Generate N diverse reformulations of the query via LLM.
        Different phrasings hit different embedding regions → broader recall.
        Always includes the original query; returns [original, v1, v2, ...].
        Falls back to [query] on failure.
        """
        messages = [
            {
                "role": "system",
                "content": (
                    f"Generate {n} diverse reformulations of the user's query to improve document retrieval. "
                    "Each reformulation should approach the topic from a different angle or use different keywords. "
                    f"Return exactly {n} queries, one per line, numbered 1. 2. etc. "
                    "No explanations — only the queries."
                ),
            },
            {"role": "user", "content": query},
        ]
        try:
            response = await asyncio.to_thread(
                self.openrouter.chat_completion,
                messages=messages,
                model=INTERNAL_LLM_MODEL,
                temperature=0.7,
                max_tokens=150,
            )
            content = response.get("content", "").strip()
            variants = []
            for line in content.split("\n"):
                line = re.sub(r"^\d+\.\s*", "", line.strip())
                if line:
                    variants.append(line)
            return [query] + variants[:n]
        except Exception as e:
            logger.warning(f"Query variant generation failed: {e}")
            return [query]

    async def _multi_query_search(
        self, bot_id: str, queries: List[str], top_k: int
    ) -> List[Dict]:
        """
        Multi-Query Fusion — step 2:
        Embed each query variant independently (HyDE already applied to first),
        run _hybrid_search (rerank=False) for each, merge via RRF, then do ONE
        final Cross-Encoder rerank pass on the merged pool.

        Skipping per-variant reranking saves ~(n_variants - 1) expensive CrossEncoder
        calls (e.g., 3 calls → 1 call) while preserving quality via final rerank.
        """
        async def _search_one(q: str) -> List[Dict]:
            embedding = await asyncio.to_thread(self._embed_with_retry, q)
            return await self._hybrid_search(bot_id, q, embedding, top_k, False)

        results_per_query = await asyncio.gather(
            *[_search_one(q) for q in queries], return_exceptions=True
        )

        k = 60  # standard RRF constant
        rrf_scores: Dict[str, float] = {}
        best_payload: Dict[str, Dict] = {}

        for results in results_per_query:
            if isinstance(results, BaseException):
                continue
            for rank, result in enumerate(results):
                key = result.get("text", "")[:200]
                rrf_scores[key] = rrf_scores.get(key, 0.0) + 1.0 / (k + rank + 1)
                if key not in best_payload or result.get("rrf_score", 0) > best_payload[key].get("rrf_score", 0):
                    best_payload[key] = result

        if not best_payload:
            return []

        merged = []
        for key, rrf_score in sorted(rrf_scores.items(), key=lambda x: x[1], reverse=True):
            item = best_payload[key].copy()
            item["hybrid_score"] = min(rrf_score * k, 1.0)  # normalized cross-query RRF
            merged.append(item)

        # ONE final rerank on the merged pool (replaces 3× per-variant reranks)
        reranker = self._get_reranker()
        if reranker and merged:
            import numpy as np
            try:
                primary_query = queries[0]
                pool = merged[: top_k * 2]
                pairs = [[primary_query, item["text"]] for item in pool]
                scores = reranker.predict(pairs, show_progress_bar=False)
                if isinstance(scores, float):
                    scores = [scores]
                norm_scores = _sigmoid(np.array(scores))
                for i, item in enumerate(pool):
                    rrf_base = item.get("hybrid_score", 0.0)
                    item["hybrid_score"] = 0.7 * float(norm_scores[i]) + 0.3 * rrf_base
                merged = sorted(pool, key=lambda x: x["hybrid_score"], reverse=True)
            except Exception as e:
                logger.warning(f"Multi-query final rerank failed: {e}")
                merged.sort(key=lambda x: x["hybrid_score"], reverse=True)
        else:
            merged.sort(key=lambda x: x["hybrid_score"], reverse=True)

        return merged[: top_k * 2]

    # ──────────────────────────────────────────────────────────────────────────
    # Contextual Retrieval (Anthropic technique)
    # ──────────────────────────────────────────────────────────────────────────

    async def _generate_contextual_prefix_batch(
        self, doc_text: str, chunk_texts: List[str], max_chunks: int = 0
    ) -> List[str]:
        """
        Contextual Retrieval — generate a 1-2 sentence situating context for each chunk.
        Batches 8 chunks per LLM call (down from 50 separate calls) and removes the
        cap so all chunks get prefixes.

        The prefix is prepended to the chunk *before* embedding so that the vector
        captures where the chunk lives within the whole document (Anthropic technique).

        Caps at max_chunks to bound indexing cost. Chunks beyond the cap get empty prefix.
        Falls back gracefully on any individual failure.
        """
        # Truncate full doc text to keep prompts manageable
        doc_summary = doc_text[:4000]
        BATCH_SIZE = 8

        async def _prefix_batch(chunk_batch: List[tuple]) -> List[str]:
            """Generate prefixes for a batch of (idx, chunk) tuples in one LLM call."""
            chunks_xml = "\n".join(
                f'<chunk id="{i}">{text[:600]}</chunk>'
                for i, text in chunk_batch
            )
            messages = [
                {
                    "role": "user",
                    "content": (
                        f"<document>\n{doc_summary}\n</document>\n\n"
                        f"{chunks_xml}\n\n"
                        "For each chunk, write 1-2 sentences that situate it within the "
                        "overall document (which section, what topic). "
                        'Reply as JSON: {"prefixes": [{"id": 0, "text": "..."}, ...]}\n'
                        "Only the JSON object, no other text."
                    ),
                }
            ]
            try:
                response = await asyncio.to_thread(
                    self.openrouter.chat_completion,
                    messages=messages,
                    model=INTERNAL_LLM_MODEL,
                    temperature=0.1,
                    max_tokens=80 * len(chunk_batch),
                )
                content = response.get("content", "").strip()
                # Try JSON parse
                import json as _json
                match = re.search(r'\{.*\}', content, re.DOTALL)
                if match:
                    parsed = _json.loads(match.group())
                    prefix_map = {p["id"]: p["text"] for p in parsed.get("prefixes", [])}
                    return [prefix_map.get(i, "") for i, _ in chunk_batch]
            except Exception:
                logger.warning("contextual_prefix_parse_failed")
            # Fallback: one-by-one for this batch
            fallback = []
            for i, text in chunk_batch:
                try:
                    resp = await asyncio.to_thread(
                        self.openrouter.chat_completion,
                        messages=[{"role": "user", "content": (
                            f"<document>\n{doc_summary}\n</document>\n\n"
                            f"<chunk>\n{text[:600]}\n</chunk>\n\n"
                            "Write 1-2 sentences that situate this chunk within the document. "
                            "Reply with ONLY the context sentences."
                        )}],
                        model=INTERNAL_LLM_MODEL,
                        temperature=0.1,
                        max_tokens=80,
                    )
                    fallback.append(resp.get("content", "").strip())
                except Exception:
                    fallback.append("")
            return fallback

        # Build batches of (index, chunk_text)
        indexed = list(enumerate(chunk_texts))
        batches = [indexed[i:i + BATCH_SIZE] for i in range(0, len(indexed), BATCH_SIZE)]

        batch_results = await asyncio.gather(
            *[_prefix_batch(b) for b in batches], return_exceptions=True
        )

        result = [""] * len(chunk_texts)
        for batch, batch_result in zip(batches, batch_results):
            if isinstance(batch_result, BaseException):
                continue
            for (idx, _), prefix in zip(batch, batch_result):
                result[idx] = prefix
        return result

    async def _get_kb_version(self, bot_id: str) -> int:
        """Knowledge-base version used to scope cached RAG/CRAG responses."""
        try:
            redis = self.cache.redis if self.cache else None
            if redis is None:
                return 0
            value = await redis.get(f"kb_version:{bot_id}")
            if value is None:
                return 0
            if isinstance(value, bytes):
                value = value.decode("utf-8")
            return int(value)
        except Exception:
            logger.warning("kb_version_read_failed", extra={"bot_id": bot_id})
            return 0

    async def _bump_kb_version(self, bot_id: str) -> int:
        """Increment bot KB version after ingest/delete/reindex."""
        try:
            redis = self.cache.redis if self.cache else None
            if redis is None:
                return 0
            return int(await redis.incr(f"kb_version:{bot_id}"))
        except Exception:
            logger.warning("kb_version_bump_failed", extra={"bot_id": bot_id})
            return 0

    @staticmethod
    def _is_response_cache_safe(
        bot_config: Dict[str, Any],
        conversation_history: Optional[List[Dict[str, str]]],
    ) -> bool:
        """Only cache stateless responses; memory/history/group context are user-scoped."""
        if conversation_history:
            return False
        if bot_config.get("enable_memory", True):
            return False
        if bot_config.get("user_id"):
            return False
        if bot_config.get("group_participants") or bot_config.get("group_recent_messages"):
            return False
        return True

    @staticmethod
    def _cache_relevant_config(bot_config: Dict[str, Any]) -> Dict[str, Any]:
        keys = {
            "model",
            "temperature",
            "max_tokens",
            "system_prompt",
            "domain",
            "top_k",
            "similarity_threshold",
            "enable_knowledge_graph",
            "_lightrag_mode",
        }
        return {key: bot_config.get(key) for key in sorted(keys) if key in bot_config}

    async def chat(
        self,
        bot_id: str,
        query: str,
        bot_config: Optional[Dict[str, Any]] = None,
        conversation_history: Optional[List[Dict[str, str]]] = None,
        session_id: Optional[str] = None,
        top_k: int = 5,
        use_cache: bool = True
    ) -> Dict[str, Any]:
        """
        Chat with bot using OpenRouter with Agentic Query Rewriting and Reranking.
        """
        start_time = time.time()
        session_id = session_id or str(uuid.uuid4())
        bot_config = bot_config or {}

        # Inject domain profile settings into bot_config (profile values are fallbacks;
        # explicit bot_config values take precedence).
        try:
            from app.services.domain_config import get_domain_profile
            domain = bot_config.get("domain", "general")
            profile = get_domain_profile(domain)
            # Only apply profile top_k if not overridden by caller arg or bot config
            effective_top_k = bot_config.get("top_k", profile.retrieval_k)
            # Propagate lightrag mode so _prepare_chat_context can use it
            bot_config.setdefault("_lightrag_mode", profile.lightrag_mode)
            # If KG is enabled but domain says don't use lightrag, honour the domain default
            # (bot_config.enable_knowledge_graph remains the final gate)
            domain_prompt_suffix = profile.system_prompt_suffix
        except Exception:
            effective_top_k = top_k
            domain_prompt_suffix = ""

        kb_version = await self._get_kb_version(bot_id)
        cache_allowed = (
            use_cache
            and self.cache
            and self._is_response_cache_safe(bot_config, conversation_history)
        )
        cache_key_data = {
            "bot_id": bot_id,
            "query": hashlib.md5(query.encode()).hexdigest(),
            "config": self._cache_relevant_config(bot_config),
            "top_k": effective_top_k,
            "kb_version": kb_version,
        }

        # 1. Check Redis cache for stateless requests only.
        if cache_allowed:
            try:
                cached = await self.cache.get("rag:chat", cache_key_data)
                if cached:
                    logger.info(f"Cache HIT for query: {query[:50]}...")
                    cached["from_cache"] = True
                    cached["response_time"] = round(time.time() - start_time, 3)
                    return cached
            except Exception as e:
                logger.warning(f"Redis cache read error: {e}")
        
        # Prepare context (retrieval, reranking, agent logs)
        prep = await self._prepare_chat_context(bot_id, query, bot_config, effective_top_k)
        search_query = prep["search_query"]
        filtered_results = prep["filtered_results"]
        agent_logs = prep["agent_logs"]
        context = prep["context"]
        sources = prep["sources"]
        crag_status = prep.get("crag_status", "relevant")

        # ── Memory: Retrieve user memories before generation ──────────────
        user_id = bot_config.get("user_id")
        enable_memory = bot_config.get("enable_memory", True)
        user_memories: list = []
        if enable_memory and user_id:
            user_memories = await memory_service.search(
                query=query,
                user_id=user_id,
                bot_id=bot_id,
                top_k=getattr(settings, "MEM0_TOP_K", 5),
            )
        # ──────────────────────────────────────────────────────────────────

        try:
            # 7. LLM Generation
            agent_logs.append({
                "step": "Answer Synthesis",
                "description": f"Generating grounded response using {len(filtered_results)} retrieved segments.",
                "timestamp": datetime.utcnow().isoformat()
            })

            base_system_prompt = bot_config.get("system_prompt", "You are a helpful assistant.")
            # Append domain-specific suffix if not already present
            if domain_prompt_suffix and domain_prompt_suffix.strip() not in base_system_prompt:
                base_system_prompt = base_system_prompt + domain_prompt_suffix
            # ── CRAG: Inject relevance signal into system prompt ─────────
            if crag_status == "no_context":
                base_system_prompt += (
                    "\n\n[CRAG SIGNAL: Knowledge base does not contain relevant information for this query. "
                    "You MUST explicitly tell the user that you don't have this information in your knowledge base "
                    "rather than guessing or fabricating an answer. Do NOT use any retrieved context below.]"
                )
            elif crag_status == "ambiguous":
                base_system_prompt += (
                    "\n\n[CRAG SIGNAL: Knowledge base only partially covers this topic. "
                    "Answer based on what is available, but clearly indicate uncertainty and "
                    "recommend the user verify with authoritative sources.]"
                )
            # ── Group chat context ────────────────────────────────────────
            group_participants = bot_config.get("group_participants")
            if group_participants:
                base_system_prompt += _build_group_context_block(bot_config)
            # ─────────────────────────────────────────────────────────────
            model = bot_config.get("model", "google/gemini-3.1-flash-lite")
            temperature = bot_config.get("temperature", 0.7)
            max_tokens = bot_config.get("max_tokens", 1000)

            if context:
                if "{{context}}" in base_system_prompt:
                    # User actively provided placeholder {{context}}
                    effective_system_prompt = base_system_prompt.replace("{{context}}", context)
                else:
                    # Fallback for old bots without placeholder (append to end)
                    effective_system_prompt = (
                        f"{base_system_prompt}\n\n"
                        "CÁCH TRẢ LỜI (QUY TẮC BẮT BUỘC):\n"
                        "1. Sử dụng DUY NHẤT ngữ cảnh được cung cấp dưới đây để trả lời.\n"
                        "2. Nếu thông tin không có trong ngữ cảnh, hãy nói rằng bạn không biết, KHÔNG tự chế câu trả lời.\n"
                        "3. TRÍCH DẪN NGUỒN: Sử dụng ký hiệu [[n]] (ví dụ [[1]], [[2]]) ngay sau câu hoặc cụm từ trích dẫn thông tin từ Segment tương ứng.\n"
                        "4. Luôn trả lời bằng ngôn ngữ của người dùng (Tiếng Việt).\n\n"
                        f"CONTEXT (TÀI LIỆU TRÍCH XUẤT):\n{context}"
                    )
                user_content = query
            else:
                # If no context is found, clear the placeholder from prompt
                effective_system_prompt = base_system_prompt.replace("{{context}}", "")
                user_content = query

            # ── Memory: Prepend user memory block to the system prompt ────────
            memory_block = memory_service.build_memory_prompt_block(user_memories)
            effective_system_prompt = memory_block + effective_system_prompt
            # ─────────────────────────────────────────────────────────────────

            messages = [
                {"role": "system", "content": effective_system_prompt},
                {"role": "user", "content": user_content}
            ]
            
            if conversation_history:
                _hist_limit = int(bot_config.get("conversation_history_limit", 5))
                messages[1:1] = conversation_history[-_hist_limit:]

            llm_response = await asyncio.to_thread(
                self._chat_with_retry,
                messages=messages,
                model=model,
                temperature=temperature,
                max_tokens=max_tokens
            )
            
            response_text = llm_response["content"]
            usage = llm_response["usage"]
            response_time = time.time() - start_time

            # 8. Reasoning Summary
            reasoning = prep["reasoning"]

            # 9. Prepare result
            result = {
                "response": response_text,
                "sources": sources,
                "retrieved_chunks": filtered_results,
                "agent_logs": agent_logs,
                "reasoning": reasoning,
                "model": model,
                "usage": {
                    "prompt_tokens": usage["prompt_tokens"],
                    "completion_tokens": usage["completion_tokens"],
                    "total_tokens": usage["total_tokens"]
                },
                "response_time": round(response_time, 3),
                "session_id": session_id,
                "from_cache": False,
                "message_id": str(uuid.uuid4()),
                "search_query": search_query,
                "system_prompt": effective_system_prompt,
                "memories_used": user_memories,
                "lightrag_entities": prep.get("lightrag_entities", []),
            }

            # ── Memory: Save this conversation turn (non-blocking) ────────
            if enable_memory and user_id:
                asyncio.create_task(memory_service.add(
                    messages=[
                        {"role": "user", "content": query},
                        {"role": "assistant", "content": response_text},
                    ],
                    user_id=user_id,
                    bot_id=bot_id,
                    session_id=session_id,
                ))
            # ─────────────────────────────────────────────────────────────

            # 9. Cache response
            if cache_allowed:
                try:
                    cache_payload = result.copy()
                    cache_payload.pop("session_id", None)
                    cache_payload.pop("message_id", None)
                    cache_payload["from_cache"] = True
                    await self.cache.set("rag:chat", cache_key_data, cache_payload, ttl=3600)
                except Exception as e:
                    logger.warning(f"Failed to write to Redis cache: {e}")

            # 10. Log conversation
            try:
                await self._log_conversation(
                    bot_id=bot_id,
                    session_id=session_id,
                    user_id=bot_config.get("user_id"),
                    user_message=query,
                    response=response_text,
                    sources=sources,
                    response_time=response_time,
                    model=model,
                    usage=usage,
                    retrieved_chunks=filtered_results,
                    reasoning=reasoning,
                    search_query=search_query,
                    agent_logs=agent_logs
                )
            except Exception as e:
                logger.error(f"Failed to log conversation: {e}")

            return result
            
        except Exception as e:
            logger.error(f"Error in chat: {str(e)}")
            raise
    
    async def _prepare_chat_context(
        self,
        bot_id: str,
        query: str,
        bot_config: Dict[str, Any],
        top_k: int = 5,
        debug_mode: bool = False
    ) -> Dict[str, Any]:
        """
        Fast retrieval pipeline — optimized for minimum latency to first token.

        Pipeline (all steps overlapped as much as possible):
          t=0:   embed(query) + rewrite(query)       ← concurrent
          t=0.4: embed done → search + lightrag start ← no waiting for rewrite
          t=1.7: search done → CRAG starts immediately ← lightrag still running
          t=1.5: rewrite done (overlaps with search)
          t=2.8: CRAG + lightrag both done            ← concurrent with each other
          t=2.8: → LLM streaming begins
        """
        import time as _time
        _t0 = _time.time()

        use_kg = bot_config.get("enable_knowledge_graph", False)
        lightrag_mode = bot_config.get("_lightrag_mode", "hybrid")
        similarity_threshold = bot_config.get("similarity_threshold", 0.15)

        agent_logs = [{
            "step": "Analyzing Query",
            "description": f"Embedding + rewriting query: '{query[:50]}...'",
            "timestamp": datetime.utcnow().isoformat()
        }]

        # ── Step 1: embed(query) + rewrite(query) — CONCURRENT ─────────────
        # Embed does NOT need the rewritten query — fire both immediately.
        logger.debug("perf_embed_rewrite_concurrent")

        async def _cached_embed(q: str) -> list:
            """Embed with Redis cache (singleflight protected)."""
            cache_key = f"q:{self.openrouter.embedding_model}:{q.strip().lower()}"
            cached = await self.cache.get("emb", cache_key)
            if cached is not None:
                logger.debug("perf_embed_cache_hit")
                return cached
            result = await asyncio.to_thread(self._embed_with_retry, q)
            await self.cache.set("emb", cache_key, result, ttl=86400)
            return result

        async def _cached_rewrite(q: str) -> str:
            """Rewrite with cache. Skip for short queries."""
            if len(q.strip()) < 20 and not q.endswith("?"):
                return q
            cache_key = f"rw:{q}"
            cached = await self.cache.get("rewrite", cache_key)
            if cached is not None:
                return cached
            result = await self._rewrite_query(q)
            if result != q:
                await self.cache.set("rewrite", cache_key, result, ttl=1800)
            return result

        embed_task = asyncio.ensure_future(_cached_embed(query))
        rewrite_task = asyncio.ensure_future(_cached_rewrite(query))

        # ── Step 2: search + lightrag — start as soon as embed is ready ────
        # DO NOT await rewrite_task here — search uses original-query embedding.
        query_embedding = await embed_task
        logger.debug("perf_embed_done", extra={"duration_s": round(_time.time()-_t0, 2)})

        agent_logs.append({
            "step": "Knowledge Retrieval",
            "description": "Hybrid search" + (" + knowledge graph" if use_kg else "") + " in parallel.",
            "timestamp": datetime.utcnow().isoformat()
        })

        async def _run_lightrag(bid: str, q: str, mode: str = "hybrid") -> str:
            try:
                from app.services.lightrag_service import get_lightrag_service
                svc = get_lightrag_service(bot_id=bid)
                result = await asyncio.wait_for(svc.query(q, mode=mode), timeout=8.0)
                logger.info(f"[PERF] KG contributed: yes, took={mode}, bot={bid}")
                return result or ""
            except asyncio.TimeoutError:
                logger.warning(f"[PERF] KG timed out (8s) for bot={bid}")
                return ""
            except Exception as exc:
                logger.warning(f"LightRAG query failed: {exc}")
                return ""

        _t1 = _time.time()
        # search uses original query embedding (not HyDE — saves 1 LLM call)
        search_task = asyncio.ensure_future(
            self._hybrid_search(bot_id, query, query_embedding, top_k)
        )
        # lightrag uses original query (rewrite not ready yet)
        lightrag_task = asyncio.ensure_future(
            _run_lightrag(bot_id, query, lightrag_mode) if use_kg
            else asyncio.sleep(0, result="")
        )

        # ── Step 3: CRAG starts as soon as search is done ──────────────────
        # lightrag keeps running concurrently — we don't block on it here.
        search_results = await search_task
        logger.debug("perf_search_done", extra={"duration_s": round(_time.time()-_t1, 2)})

        if isinstance(search_results, BaseException):
            logger.warning(f"Hybrid search failed: {search_results}")
            search_results = []

        similarity_threshold = bot_config.get("similarity_threshold", 0.15)
        filtered_results = [
            r for r in (search_results or [])
            if r.get("hybrid_score", 0) >= similarity_threshold
        ]
        if not filtered_results and search_results:
            filtered_results = search_results

        if self.reranker:
            agent_logs.append({
                "step": "Cross-Encoder Reranking",
                "description": f"Refining {len(search_results)} candidates for maximum relevance.",
                "status": "done",
                "timestamp": datetime.utcnow().isoformat()
            })

        # Rewrite should be done by now (started at t=0, takes ~1.5s; search took ~1.3s from t=0.4)
        search_query = await rewrite_task

        # CRAG + wait for lightrag — run concurrently
        _t2 = _time.time()

        async def _cached_crag(q: str, chunks: list) -> str:
            """CRAG verdict cache keyed on query + top chunk IDs."""
            chunk_ids = hashlib.sha256(
                "|".join(c.get("text", "")[:100] for c in chunks[:5]).encode()
            ).hexdigest()[:16]
            cache_key = {
                "bot_id": bot_id,
                "kb_version": await self._get_kb_version(bot_id),
                "chunk_ids": chunk_ids,
                "query": q[:200],
            }
            cached = await self.cache.get("crag", cache_key)
            if cached is not None:
                return cached
            verdict = await self._crag_classify(q, chunks)
            await self.cache.set("crag", cache_key, verdict, ttl=3600)
            return verdict

        crag_task = asyncio.ensure_future(_cached_crag(search_query, filtered_results))

        # ── Non-blocking KG: Don't wait for LightRAG in gather ──────────────
        # CRAG is essential for quality → await it. LightRAG runs concurrently
        # with a short grace period after CRAG completes. If KG isn't ready
        # within 3s after CRAG, we proceed without it (KG adds value but is not
        # required for a correct answer). Old behaviour: blocked 10s every time.
        crag_status = await crag_task
        _kg_deadline = _time.time() + 3.0
        try:
            lightrag_raw = await asyncio.wait_for(
                lightrag_task,
                timeout=max(0.1, _kg_deadline - _time.time()),
            )
        except (asyncio.TimeoutError, asyncio.CancelledError):
            logger.info(f"[PERF] KG skipped — not ready in time for bot={bot_id}")
            lightrag_raw = ""

        if isinstance(crag_status, BaseException):
            logger.warning(f"CRAG failed: {crag_status}")
            crag_status = "ambiguous"
        if isinstance(lightrag_raw, BaseException):
            logger.warning(f"LightRAG gather error: {lightrag_raw}")
            lightrag_raw = ""

        logger.debug("perf_crag_lightrag_done", extra={"phase_s": round(_time.time()-_t2, 2), "total_prep_s": round(_time.time()-_t0, 2)})
        agent_logs.append({
            "step": "CRAG Relevance Check",
            "description": (
                "Knowledge base contains a relevant answer." if crag_status == "relevant"
                else "Knowledge base partially matches — answering with caution." if crag_status == "ambiguous"
                else "Knowledge base does not contain relevant information for this query."
            ),
            "status": crag_status,
            "timestamp": datetime.utcnow().isoformat(),
        })
        # ─────────────────────────────────────────────────────────────────

        context_docs = []
        sources = []

        # Smart highlights — only for top 5 (highlight extraction is regex-heavy, skip on tail)
        for idx, result in enumerate(filtered_results):
            doc_id = idx + 1
            result["highlights"] = (
                self._extract_smart_highlights(search_query, result.get("text", ""))
                if idx < 5 else []
            )

            # Parent-Child: use the richer parent chunk for LLM context; child was used for matching
            display_text = result.get("parent_text") or result["text"]
            context_docs.append(f"[[{doc_id}]] Source: {result.get('source', 'Unknown')}\n{display_text}")
            source_name = result.get("source", "Unknown")
            if source_name not in sources:
                sources.append(source_name)

        # --- LightRAG context post-processing ---
        lightrag_context = ""
        lightrag_entities: list = []
        if lightrag_raw and not lightrag_raw.startswith("Error") and "Sorry, I'm not able to provide an answer" not in lightrag_raw:
            lightrag_context = f"\n\n--- CONTEXT TỪ KNOWLEDGE GRAPH (ENTITIES & RELATIONSHIPS) ---\n{lightrag_raw}\n(Dùng context trên để bổ sung vào câu trả lời, không trích dẫn nguyên văn)\n\n"
            # Extract entity names so frontend can highlight them in the knowledge graph
            lightrag_entities = _extract_lightrag_entity_names(lightrag_raw)
            agent_logs.append({
                "step": "Knowledge Graph Retrieval",
                "description": f"Successfully traversed the entity-relationship graph. Found {len(lightrag_entities)} relevant entities.",
                "timestamp": datetime.utcnow().isoformat()
            })
        
        context = "\n---\n".join(context_docs)
        if lightrag_context:
            context += lightrag_context
        reasoning = f"I've analyzed {len(filtered_results)} segments from {len(sources)} documents. "
        if filtered_results:
            top_score = filtered_results[0].get('hybrid_score', 0)
            reasoning += f"The most relevant source is '{filtered_results[0].get('source')}' with a confidence score of {top_score:.2f}."

        return {
            "search_query": search_query,
            "filtered_results": filtered_results,
            "agent_logs": agent_logs,
            "context": context,
            "sources": sources,
            "reasoning": reasoning,
            "lightrag_entities": lightrag_entities,
            "crag_status": crag_status,
            "hyde_hypothesis": "",
            "multi_query_variants": [],
        }

    def _extract_smart_highlights(self, query: str, text: str) -> List[str]:
        """Backend smart highlight logic: filters noise and identifies key terms."""
        if not query or not text:
            return []
            
        # 1. Clean query and extract candidates
        # Remove common punctuation and split
        words = re.findall(r'\w+', query.lower())
        
        # 2. Vietnamese stop words + common conversational filler
        stop_words = {
            'tàu', 'xe', 'chạy', 'cho', 'của', 'là', 'có', 'không', 'trong', 'về', 'nhé', 
            'em', 'được', 'tại', 'vào', 'ngày', 'lúc', 'để', 'đến', 'với', 'cần', 'hỏi'
        }
        
        candidates = [w for w in words if len(w) > 1 and w not in stop_words]
        
        # 3. Look for "High Value" patterns (Mã hiệu, số hiệu like SNT1, SE4, 24/11)
        # We can use regex to find codes in the text that are also partially in the query
        highlights = []
        for cand in candidates:
            # Simple match for now, could be improved with fuzzy match or NER
            if cand in text.lower():
                # Try to find the original casing from the text
                matches = re.findall(re.escape(cand), text, re.IGNORECASE)
                if matches:
                    highlights.extend(list(set(matches)))
        
        # Unique and sorted by length
        return sorted(list(set(highlights)), key=len, reverse=True)[:10]

    async def chat_stream(
        self,
        bot_id: str,
        query: str,
        bot_config: Optional[Dict[str, Any]] = None,
        conversation_history: Optional[List[Dict[str, str]]] = None,
        session_id: Optional[str] = None,
        top_k: int = 5,
    ) -> AsyncGenerator[Dict[str, Any], None]:
        """Streaming version of chat using AsyncGenerator."""
        logger.info(f"[CHAT_STREAM] Started for bot={bot_id}, query={query[:50]}")
        start_time = time.time()
        session_id = session_id or str(uuid.uuid4())
        bot_config = bot_config or {}

        # Inject domain profile settings (same logic as chat())
        try:
            from app.services.domain_config import get_domain_profile
            domain = bot_config.get("domain", "general")
            profile = get_domain_profile(domain)
            effective_top_k = bot_config.get("top_k", profile.retrieval_k)
            bot_config.setdefault("_lightrag_mode", profile.lightrag_mode)
            domain_prompt_suffix = profile.system_prompt_suffix
        except Exception:
            effective_top_k = top_k
            domain_prompt_suffix = ""

        # ── Memory search + RAG prep — run concurrently ──────────────────
        user_id = bot_config.get("user_id")
        enable_memory = bot_config.get("enable_memory", True)

        async def _fetch_memories():
            if not (enable_memory and user_id):
                return []
            return await memory_service.search(
                query=query, user_id=user_id, bot_id=bot_id,
                top_k=getattr(settings, "MEM0_TOP_K", 5),
            )

        prep, user_memories = await asyncio.gather(
            self._prepare_chat_context(bot_id, query, bot_config, effective_top_k),
            _fetch_memories(),
        )
        # ─────────────────────────────────────────────────────────────────

        search_query = prep["search_query"]
        filtered_results = prep["filtered_results"]
        agent_logs = prep["agent_logs"]
        context = prep["context"]
        sources = prep["sources"]
        reasoning = prep["reasoning"]
        lightrag_entities = prep.get("lightrag_entities", [])
        crag_status = prep.get("crag_status", "relevant")

        # Yield metadata (FINAL context results)
        logger.info(f"[CHAT_STREAM] Yielding metadata, filtered_results count={len(filtered_results)}")
        yield {
            "type": "metadata",
            "sources": sources,
            "retrieved_chunks": filtered_results,
            "agent_logs": agent_logs,
            "reasoning": reasoning,
            "search_query": search_query,
            "session_id": session_id,
            "lightrag_entities": lightrag_entities,
        }
        logger.info(f"[CHAT_STREAM] Metadata yielded, proceeding to LLM generation")

        # 5. Answer Synthesis Log
        agent_logs.append({
            "step": "Answer Synthesis",
            "description": f"Generating grounded response using {len(filtered_results)} retrieved segments.",
            "timestamp": datetime.utcnow().isoformat()
        })

        # Setup LLM Generation
        base_system_prompt = bot_config.get("system_prompt", "You are a helpful assistant.")
        if domain_prompt_suffix and domain_prompt_suffix.strip() not in base_system_prompt:
            base_system_prompt = base_system_prompt + domain_prompt_suffix
        # ── CRAG: Inject relevance signal ────────────────────────────────
        if crag_status == "no_context":
            base_system_prompt += (
                "\n\n[CRAG SIGNAL: Knowledge base does not contain relevant information for this query. "
                "You MUST explicitly tell the user that you don't have this information in your knowledge base "
                "rather than guessing or fabricating an answer. Do NOT use any retrieved context below.]"
            )
        elif crag_status == "ambiguous":
            base_system_prompt += (
                "\n\n[CRAG SIGNAL: Knowledge base only partially covers this topic. "
                "Answer based on what is available, but clearly indicate uncertainty and "
                "recommend the user verify with authoritative sources.]"
            )
        # ── Group chat context ────────────────────────────────────────────
        group_participants = bot_config.get("group_participants")
        if group_participants:
            base_system_prompt += _build_group_context_block(bot_config)
        # ─────────────────────────────────────────────────────────────────
        model = bot_config.get("model", "google/gemini-3.1-flash-lite")
        temperature = bot_config.get("temperature", 0.7)
        max_tokens = bot_config.get("max_tokens", 1000)

        if context:
            effective_system_prompt = (
                f"{base_system_prompt}\n\n"
                "CÁCH TRẢ LỜI (QUY TẮC BẮT BUỘC):\n"
                "1. Sử dụng DUY NHẤT ngữ cảnh được cung cấp dưới đây để trả lời.\n"
                "2. Nếu thông tin không có trong ngữ cảnh, hãy nói rằng bạn không biết, KHÔNG tự chế câu trả lời.\n"
                "3. TRÍCH DẪN NGUỒN: Sử dụng ký hiệu [[n]] (ví dụ [[1]], [[2]]) ngay sau câu hoặc cụm từ trích dẫn thông tin từ Segment tương ứng.\n"
                "4. Luôn trả lời bằng ngôn ngữ của người dùng (Tiếng Việt).\n\n"
                f"CONTEXT (TÀI LIỆU TRÍCH XUẤT):\n{context}"
            )
            user_content = query
        else:
            effective_system_prompt = base_system_prompt
            user_content = query

        # ── Memory: Prepend user memory block to the system prompt ────────
        memory_block = memory_service.build_memory_prompt_block(user_memories)
        if memory_block:
            effective_system_prompt = memory_block + effective_system_prompt
        # ─────────────────────────────────────────────────────────────

        messages = [
            {"role": "system", "content": effective_system_prompt},
            {"role": "user", "content": user_content}
        ]
        if conversation_history:
            messages[1:1] = conversation_history[-5:]

        # Streaming from OpenRouter
        full_response = ""
        try:
            chunk_count = 0
            async for chk in self.openrouter.chat_completion_stream(
                messages=messages,
                model=model,
                temperature=temperature,
                max_tokens=max_tokens,
            ):
                if getattr(chk, "choices", None) and len(chk.choices) > 0:
                    delta = getattr(chk.choices[0], "delta", None)
                    if delta and getattr(delta, "content", None):
                        content = delta.content
                        full_response += content
                        chunk_count += 1
                        yield {"type": "content", "content": content}

            logger.info(f"[CHAT_STREAM] Finished streaming {chunk_count} chunks, response length={len(full_response)}")
            if len(full_response) == 0:
                logger.warning("Warning: Generated response is empty!")

            # Log conversation BEFORE yielding "done" — this is critical!
            # The generator may be closed by StreamingResponse after the last yield,
            # so we must do all post-processing BEFORE that final yield.
            try:
                logger.info(f"[STREAM] Logging conversation for session={session_id}")
                await self._log_conversation(
                    bot_id=bot_id,
                    session_id=session_id or "default",
                    user_id=bot_config.get("user_id"),
                    user_message=query,
                    response=full_response,
                    sources=sources,
                    response_time=time.time() - start_time,
                    model=model,
                    usage={"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0},
                    retrieved_chunks=filtered_results,
                    reasoning=reasoning,
                    search_query=search_query,
                    agent_logs=agent_logs
                )
                logger.info(f"[STREAM] Successfully logged conversation for session={session_id}")
            except Exception as log_err:
                logger.error(f"[STREAM] Failed to log conversation: {log_err}", exc_info=True)

            # Memory saving (non-blocking)
            if enable_memory and user_id and full_response:
                asyncio.create_task(memory_service.add(
                    messages=[
                        {"role": "user", "content": query},
                        {"role": "assistant", "content": full_response},
                    ],
                    user_id=user_id,
                    bot_id=bot_id,
                ))

            yield {"type": "done"}
            
        except Exception as e:
            logger.error(f"Error in chat_stream: {e}")
            yield {"type": "error", "message": str(e)}

    async def _log_conversation(
        self,
        bot_id: str,
        session_id: str,
        user_id: Optional[str],
        user_message: str,
        response: str,
        sources: List[str],
        response_time: float,
        model: str,
        usage: Dict[str, int],
        retrieved_chunks: Optional[List[Dict]] = None,
        reasoning: Optional[str] = None,
        search_query: Optional[str] = None,
        agent_logs: Optional[List[Dict]] = None
    ):
        """Log conversation to MongoDB and maintain session list."""
        try:
            mongo_db = await get_mongodb()
            conversations_collection = mongo_db.conversations
            sessions_collection = mongo_db.sessions
            
            now = datetime.utcnow()
            
            # 1. Log the individual message
            conversation_doc = {
                "bot_id": bot_id,
                "session_id": session_id,
                "user_id": user_id,
                "user_message": user_message,
                "response": response,
                "sources": sources,
                "response_time": round(response_time, 2),
                "model": model,
                "usage": usage,
                "retrieved_chunks": retrieved_chunks,
                "reasoning": reasoning,
                "search_query": search_query,
                "agent_logs": agent_logs,
                "timestamp": now,
                "provider": "openrouter"
            }
            
            await conversations_collection.insert_one(conversation_doc)
            
            # 2. Upsert session (atomic operation to prevent duplicate sessions)
            title = user_message[:50] + ("..." if len(user_message) > 50 else "")
            
            result = await sessions_collection.update_one(
                {"session_id": session_id, "bot_id": bot_id},
                {
                    "$set": {"updated_at": now},
                    "$setOnInsert": {
                        "session_id": session_id,
                        "bot_id": bot_id,
                        "user_id": user_id,
                        "title": title,
                        "created_at": now,
                    }
                },
                upsert=True
            )
            
            if result.upserted_id:
                # This was a new session — generate a better title
                logger.info(f"Created new session: {session_id} with title: {title}")
                asyncio.create_task(self._summarize_session_title(session_id, user_message))
            
            logger.debug(f"Conversation logged: session_id={session_id}")
            
        except Exception as e:
            logger.error(f"Error logging conversation: {e}", exc_info=True)

    async def _summarize_session_title(self, session_id: str, first_message: str):
        """Generate a short (3-5 words) title from the first message in Vietnamese."""
        try:
            prompt = [
                {"role": "system", "content": "You are a helpful assistant. Create a very short title (max 5 words) for a chat conversation starting with the user message below. Output ONLY the title, no quotes or prefix. The title MUST be in Vietnamese (Tiếng Việt)."},
                {"role": "user", "content": first_message}
            ]
            response = await asyncio.to_thread(
                self.openrouter.chat_completion,
                messages=prompt,
                model=INTERNAL_LLM_MODEL,
                temperature=0.3,
                max_tokens=20
            )
            title = response.get("content", "").strip().strip('"').strip("'")
            if title:
                mongo_db = await get_mongodb()
                await mongo_db.sessions.update_one(
                    {"session_id": session_id},
                    {"$set": {"title": title}}
                )
                logger.debug(f"Updated session {session_id} title to: {title}")
        except Exception as e:
            logger.warning(f"Failed to summarize session title: {e}")

    async def get_sessions(
        self,
        bot_id: str,
        user_id: Optional[str] = None,
        limit: int = 50
    ) -> List[Dict[str, Any]]:
        """List all chat sessions for a bot and user."""
        try:
            mongo_db = await get_mongodb()
            sessions_collection = mongo_db.sessions
            
            query = {"bot_id": bot_id}
            if user_id:
                query["user_id"] = user_id
            
            cursor = sessions_collection.find(query).sort("updated_at", -1).limit(limit)
            sessions = await cursor.to_list(length=limit)
            
            # Format for frontend
            return [{
                "id": s["session_id"],
                "title": s["title"],
                "updated_at": s["updated_at"].isoformat() if isinstance(s["updated_at"], datetime) else s.get("updated_at"),
                "created_at": s["created_at"].isoformat() if isinstance(s["created_at"], datetime) else s.get("created_at")
            } for s in sessions]
        except Exception as e:
            logger.error(f"Error retrieving sessions: {e}")
            return []

    async def get_chat_history(
        self,
        bot_id: str,
        user_id: Optional[str] = None,
        session_id: Optional[str] = None,
        limit: int = 50
    ) -> List[Dict[str, Any]]:
        """Retrieve recent chat history for a bot, specific session preferred."""
        try:
            mongo_db = await get_mongodb()
            conversations_collection = mongo_db.conversations
            
            query = {"bot_id": bot_id}
            if session_id:
                query["session_id"] = session_id
            elif user_id:
                query["user_id"] = user_id
            
            cursor = conversations_collection.find(query).sort("timestamp", -1).limit(limit)
            conversations = await cursor.to_list(length=limit)
            
            history = []
            for conv in reversed(conversations):
                history.append({
                    "id": f"{conv['_id']}_u",
                    "role": "user",
                    "content": conv["user_message"],
                    "timestamp": conv["timestamp"].isoformat() if isinstance(conv["timestamp"], datetime) else conv.get("timestamp")
                })
                history.append({
                    "id": f"{conv['_id']}_a",
                    "message_id": str(conv["_id"]),
                    "role": "assistant",
                    "content": conv["response"],
                    "timestamp": conv["timestamp"].isoformat() if isinstance(conv["timestamp"], datetime) else conv.get("timestamp"),
                    "sources": conv.get("sources", []),
                    "retrieved_chunks": conv.get("retrieved_chunks", []),
                    "reasoning": conv.get("reasoning", ""),
                    "agent_logs": conv.get("agent_logs", []),
                    "search_query": conv.get("search_query", "")
                })
            
            return history
        except Exception as e:
            logger.error(f"Error retrieving chat history: {e}")
            return []

    async def delete_session(
        self,
        bot_id: str,
        session_id: str,
        user_id: Optional[str] = None
    ) -> bool:
        """Delete a specific chat session and all its messages."""
        try:
            mongo_db = await get_mongodb()
            conversations_collection = mongo_db.conversations
            sessions_collection = mongo_db.sessions
            
            # Build query
            query = {"bot_id": bot_id, "session_id": session_id}
            if user_id:
                query["user_id"] = user_id
            
            # Delete messages
            conv_result = await conversations_collection.delete_many(query)
            
            # Delete session
            session_result = await sessions_collection.delete_one(query)
            
            logger.info(f"Deleted session {session_id} and {conv_result.deleted_count} messages")
            return session_result.deleted_count > 0
        except Exception as e:
            logger.error(f"Error deleting session: {e}")
            return False

    async def clear_all_history(
        self,
        bot_id: str,
        user_id: Optional[str] = None
    ) -> bool:
        """Delete all chat sessions and messages for a bot and user."""
        try:
            mongo_db = await get_mongodb()
            conversations_collection = mongo_db.conversations
            sessions_collection = mongo_db.sessions
            
            query = {"bot_id": bot_id}
            if user_id:
                query["user_id"] = user_id
            
            # Delete all messages
            conv_result = await conversations_collection.delete_many(query)
            
            # Delete all sessions
            session_result = await sessions_collection.delete_many(query)
            
            logger.info(f"Cleared all history for bot {bot_id}: {session_result.deleted_count} sessions, {conv_result.deleted_count} messages")
            return True
        except Exception as e:
            logger.error(f"Error clearing history: {e}")
            return False
    
    async def invalidate_bot_cache_async(self, bot_id: str):
        """Invalidate bot-scoped RAG caches by bumping the KB version key."""
        try:
            version = await self._bump_kb_version(bot_id)
            logger.info("bot_cache_invalidated", extra={"bot_id": bot_id, "kb_version": version})
        except Exception as e:
            logger.error("cache_invalidation_failed", extra={"bot_id": bot_id, "error": str(e)})

    def invalidate_bot_cache(self, bot_id: str):
        """Sync wrapper — schedule async invalidation as a task."""
        try:
            loop = asyncio.get_event_loop()
            if loop.is_running():
                asyncio.ensure_future(self.invalidate_bot_cache_async(bot_id))
            else:
                loop.run_until_complete(self.invalidate_bot_cache_async(bot_id))
        except RuntimeError:
            asyncio.run(self.invalidate_bot_cache_async(bot_id))


# Singleton instance
_openrouter_rag_service = None


def get_openrouter_rag_service() -> OpenRouterRAGService:
    """Get singleton instance of OpenRouterRAGService."""
    global _openrouter_rag_service
    if _openrouter_rag_service is None:
        _openrouter_rag_service = OpenRouterRAGService()
    return _openrouter_rag_service
